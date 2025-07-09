#!/usr/bin/env python
# -*- coding: utf-8 -*-

import asyncio
import logging
import argparse
from pathlib import Path
import sys
import os
import json
import tempfile
import shutil
from typing import Dict, Any, List, Optional
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Core System Setup ---
# Add the project root to the Python path to ensure all modules are found.
# This is necessary because this script is designed to be a self-contained example.
SCRIPT_DIR = Path(__file__).resolve().parent
# Assuming the script is in 'src/examples', the project root is two levels up.
# Adjust if the script's location is different relative to the 'autobyteus' and 'llm' packages.
PROJECT_ROOT = SCRIPT_DIR.parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

# --- Environment and Logging Configuration ---
# Load environment variables from a .env file if it exists.
# This is where API keys (GEMINI_API_KEY, OPENAI_API_KEY) should be stored.
try:
    from dotenv import load_dotenv
    # Assuming .env file is in the project root
    env_file_path = PROJECT_ROOT / ".env"
    if env_file_path.exists():
        print(f"Loading environment variables from: {env_file_path}")
        load_dotenv(dotenv_path=env_file_path)
    else:
        print(f"Info: .env file not found at {env_file_path}. Relying on system environment variables.")
except ImportError:
    print("Warning: python-dotenv not installed. Relying on system environment variables.")
    pass

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("PowerPointAgentApp")

# --- Autobyteus Framework Imports ---
# Import necessary components from the provided autobyteus framework files.
from autobyteus.agent.context.agent_config import AgentConfig
from autobyteus.agent.factory.agent_factory import AgentFactory
from autobyteus.agent.message.agent_input_user_message import AgentInputUserMessage
from autobyteus.agent.utils.wait_for_idle import wait_for_agent_to_be_idle
from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.llm.utils.llm_config import LLMConfig
from autobyteus.tools.functional_tool import tool

# --- Global State and Setup ---
# Create a temporary directory for storing generated images.
TEMP_DIR = Path(tempfile.mkdtemp(prefix="pptx_agent_"))
logger.info(f"Temporary directory for images created at: {TEMP_DIR}")

# --- Tool Definitions ---
# The agent will use these tools to perform its tasks. The `@tool` decorator
# automatically registers them with the framework and generates a schema for the LLM.

@tool(name="generate_slide_content")
async def generate_slide_content(user_input: str, num_slides: int = 5) -> str:
    """
    Generates structured JSON content for a presentation based on a user's topic.

    Args:
        user_input: The central topic or title for the presentation.
        num_slides: The desired number of content slides, in addition to the title slide.

    Returns:
        A JSON string containing the presentation title and a list of slide objects.
        Each slide object has a number, type, title, subtitle, bullets, notes, and an image_prompt.
    """
    logger.info(f"Generating slide content for topic: '{user_input}' ({num_slides} slides)")
    llm = LLMFactory.create_llm("GEMINI_2_5_PRO_API", llm_config=LLMConfig(temperature=0.7))
    prompt = f"""
    Based on the user input "{user_input}", generate the content for a {num_slides}-slide presentation.
    The output must be a single, valid JSON object. Do not include any text or formatting outside of the JSON.
    The JSON structure should be:
    {{
      "presentation_title": "A concise, catchy title for the presentation",
      "slides": [
        {{
          "number": 1,
          "type": "title",
          "title": "The main title of the presentation",
          "subtitle": "A brief, engaging subtitle",
          "notes": "Speaker notes for the title slide.",
          "image_prompt": "A simple, professional DALL-E 3 style prompt for a background image for the title slide."
        }},
        // ... list of {num_slides} more slides of type 'content' ...
        {{
          "number": 2,
          "type": "content",
          "title": "Title for Slide 2",
          "bullets": [
            "A key point or bullet for this slide.",
            "Another key point.",
            "A final key point, keep it concise."
          ],
          "notes": "Speaker notes explaining the content of slide 2 in more detail.",
          "image_prompt": "A simple, professional DALL-E 3 style prompt for an image that visually represents the content of this slide."
        }}
      ]
    }}
    Ensure every slide has a unique `image_prompt`. The prompts should be creative and relevant.
    """
    try:
        response = await llm.send_user_message(prompt)
        # Clean the response to extract only the JSON part
        json_text = response.content.strip().replace("```json", "").replace("```", "").strip()
        # Validate and return the JSON
        json.loads(json_text) # This will raise an error if the JSON is invalid
        logger.info(f"Successfully generated JSON content for {num_slides + 1} slides.")
        return json_text
    except Exception as e:
        logger.error(f"Failed to generate or parse slide content JSON: {e}", exc_info=True)
        # Return a fallback JSON structure on error
        fallback_json = json.dumps({
            "presentation_title": user_input,
            "slides": [{
                "number": 1, "type": "title", "title": user_input,
                "subtitle": "AI-Generated Presentation", "notes": "This is a fallback slide due to an error.",
                "image_prompt": "A generic professional background image for a business presentation."
            }]
        })
        return fallback_json
    finally:
        await llm.cleanup()

@tool(name="generate_slide_image")
async def generate_slide_image(image_prompt: str, slide_number: int) -> str:
    """
    Generates an image based on a text prompt and saves it locally.

    Args:
        image_prompt: A descriptive prompt for the image to be generated.
        slide_number: The number of the slide this image is for, used for naming the file.

    Returns:
        The local file path of the saved image.
    """
    logger.info(f"Generating image for slide {slide_number} with prompt: '{image_prompt}'")
    llm = LLMFactory.create_llm("GEMINI_2_0_FLASH_API")
    filepath = TEMP_DIR / f"slide_{slide_number}.png"
    try:
        # Note: Gemini Flash does not have a dedicated image generation modality like this.
        # This is a conceptual adaptation. For a real implementation with another model (like DALL-E),
        # the API call would be different. Here, we simulate by asking for an image and handling its data.
        # For actual Gemini image generation, you would use a specific model and API for that.
        # Since the provided files don't include an image generation model, we'll create a placeholder.
        
        # --- Placeholder Image Generation ---
        # This part replaces a real image generation call for demonstration purposes.
        img = Image.new('RGB', (1024, 768), color = (73, 109, 137))
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        try:
            # Use a common font if available, otherwise default
            font = ImageFont.truetype("arial.ttf", 40)
        except IOError:
            font = ImageFont.load_default()
        text = f"Image for Slide {slide_number}\n'{image_prompt[:50]}...'"
        draw.text((50, 50), text, font=font, fill=(255, 255, 255))
        img.save(filepath)
        logger.info(f"Generated placeholder image and saved to: {filepath}")
        return str(filepath.resolve())

    except Exception as e:
        logger.error(f"Image generation failed for slide {slide_number}: {e}", exc_info=True)
        # Create a fallback placeholder image on error
        img = Image.new('RGB', (1024, 768), 'lightgray')
        img.save(filepath)
        return str(filepath.resolve())
    finally:
        await llm.cleanup()

@tool(name="create_powerpoint_presentation")
async def create_powerpoint_presentation(slide_content_json: str, image_paths_json: str, output_path: str) -> str:
    """
    Creates a PowerPoint .pptx file from structured content and image paths.

    Args:
        slide_content_json: A JSON string containing the presentation title and slide data.
        image_paths_json: A JSON string mapping slide numbers to their local image file paths.
        output_path: The desired file path for the final .pptx presentation.

    Returns:
        A success message with the final path of the created presentation.
    """
    logger.info(f"Creating PowerPoint presentation at: {output_path}")
    try:
        content = json.loads(slide_content_json)
        image_paths = {int(k): v for k, v in json.loads(image_paths_json).items()}
    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON format for content or image paths. Details: {e}"

    prs = Presentation()
    # Set slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for slide_data in content['slides']:
        slide_num = slide_data['number']
        
        if slide_data['type'] == 'title':
            slide_layout = prs.slide_layouts[0] # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get('title', 'Presentation Title')
            if slide.placeholders[1]:
                 slide.placeholders[1].text = slide_data.get('subtitle', '')
        else: # 'content' type
            slide_layout = prs.slide_layouts[5] # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1.2))
            title_shape.text_frame.text = slide_data.get('title', '')
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            
            has_image = slide_num in image_paths and Path(image_paths[slide_num]).exists()
            
            # Add bullets
            if slide_data.get('bullets'):
                text_width = Inches(7.5) if has_image else Inches(15)
                body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), text_width, Inches(6.5))
                tf = body_shape.text_frame
                tf.clear() # Clear existing paragraph
                for i, bullet_text in enumerate(slide_data['bullets']):
                    p = tf.add_paragraph()
                    p.text = bullet_text
                    p.font.size = Pt(28)
                    p.level = 0

            # Add image
            if has_image:
                img_path = image_paths[slide_num]
                pic_left = Inches(8)
                pic_top = Inches(1.5)
                pic_width = Inches(7.5)
                pic_height = Inches(6.5) # Maintain aspect ratio
                slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_width)

        # Add speaker notes
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']

    # Ensure output directory exists and save the presentation
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info(f"Presentation saved successfully to {output_path}")
    return f"Presentation created and saved to {output_path}"


# --- Agent System Prompt ---
# This prompt guides the agent on how to use the tools in a sequence.
SYSTEM_PROMPT = """
You are a highly efficient assistant that creates PowerPoint presentations.
Your goal is to take a user's topic and create a complete .pptx file.
You must follow this workflow precisely:

1.  **Generate Content:** Call the `generate_slide_content` tool using the user's input topic. This will return a JSON string with all the text, notes, and image prompts for the presentation.

2.  **Generate Images (Loop):**
    -   Parse the JSON output from the previous step.
    -   Iterate through each slide object in the "slides" array.
    -   For **each slide**, call the `generate_slide_image` tool. Use the `image_prompt` and `number` from the current slide object as arguments.
    -   You will get a local file path for each generated image.

3.  **Collect Image Paths:**
    -   After the loop, you must create a single JSON object that maps each slide number (as a string key) to its corresponding image file path (as a string value).
    -   Example format: `{"1": "/path/to/slide_1.png", "2": "/path/to/slide_2.png", ...}`

4.  **Create Presentation:**
    -   Call the `create_powerpoint_presentation` tool.
    -   For the `slide_content_json` argument, pass the complete, original JSON string from Step 1.
    -   For the `image_paths_json` argument, pass the JSON string of the image path mapping you created in Step 3.
    -   For the `output_path` argument, use the file path provided by the user.

Do not ask for confirmation. Execute the entire workflow automatically.
"""

# --- Main Execution Block ---
async def main():
    """Main function to set up and run the agent."""
    parser = argparse.ArgumentParser(description="PowerPoint Generation Agent using AutoByteus")
    parser.add_argument("input", default="AI in Healthcare", help="The topic, title, or paragraph for the presentation.")
    parser.add_argument("--output", default="./presentation.pptx", help="Output path for the .pptx file.")
    parser.add_argument("--slides", type=int, default=3, help="Number of content slides to generate (total will be this + 1 for the title).")
    parser.add_argument("--agent-llm", default="GEMINI_2_5_PRO_API", help="LLM for agent orchestration (e.g., 'GPT_4o_API', 'GEMINI_1_5_PRO_API').")
    args = parser.parse_args()

    try:
        # Re-initialize the factory to discover all models, including from environment variables
        LLMFactory.reinitialize()
        
        logger.info("Initializing PowerPoint Generation Agent...")
        agent_llm = LLMFactory.create_llm(
            model_identifier=args.agent_llm,
            llm_config=LLMConfig(temperature=0.2)
        )

        tools = [generate_slide_content, generate_slide_image, create_powerpoint_presentation]

        agent_config = AgentConfig(
            name="PowerPointAgent",
            role="An autonomous presentation creator",
            description="Generates complete PowerPoint presentations from a topic, including text and images.",
            llm_instance=agent_llm,
            system_prompt=SYSTEM_PROMPT,
            tools=tools,
            auto_execute_tools=True,  # The agent will execute tools without asking for permission
            use_xml_tool_format=True  # Use XML format for tool calling instead of JSON
        )

        agent_factory = AgentFactory()
        agent = agent_factory.create_agent(config=agent_config)

        # The agent expects a structured message with all the required initial parameters.
        initial_message_content = f"""
        Please create a presentation based on the following input.
        User Topic: "{args.input}"
        Number of content slides: {args.slides}
        Final output file path: {Path(args.output).resolve()}
        """
        
        agent_input = AgentInputUserMessage(content=initial_message_content)

        # Start the agent and send the initial message
        agent.start()

        # Wait for the agent to be fully ready before sending message
        logger.info("Waiting for agent to be ready...")
        await wait_for_agent_to_be_idle(agent, timeout=30.0)  # Short timeout just to wait for ready state

        await agent.post_user_message(agent_input)

        # Give the agent a moment to process the message before checking if it's idle
        await asyncio.sleep(1.0)

        # Wait for the agent to finish its entire workflow and become idle.
        # Timeout is set to 10 minutes to allow for LLM and image generation latencies.
        logger.info("Agent is running... Waiting for the presentation to be completed.")
        await wait_for_agent_to_be_idle(agent, timeout=600.0)

        logger.info("✅ Agent has completed its task.")
        final_output_path = Path(args.output).resolve()
        logger.info(f"Presentation should be available at: {final_output_path}")

        # Copy generated images to the same directory as the presentation for reference
        output_dir = final_output_path.parent
        images_dir = output_dir / f"{final_output_path.stem}_images"
        if TEMP_DIR.exists() and any(TEMP_DIR.iterdir()):
            shutil.copytree(TEMP_DIR, images_dir, dirs_exist_ok=True)
            logger.info(f"📁 Supporting images have been copied to: {images_dir}")
        
    except Exception as e:
        logger.error(f"An error occurred during the agent execution: {e}", exc_info=True)
    finally:
        # Cleanup
        if 'agent' in locals() and agent.is_running:
            await agent.stop()
            logger.info("Agent stopped.")
        
        # Remove the temporary directory
        shutil.rmtree(TEMP_DIR, ignore_errors=True)
        logger.info(f"Temporary directory {TEMP_DIR} cleaned up.")


if __name__ == "__main__":
    # Ensure API keys are set
    if not os.getenv("GEMINI_API_KEY") or not os.getenv("OPENAI_API_KEY"):
         logger.error("API keys for Gemini and/or OpenAI are not set in the environment.")
         logger.error("Please set GEMINI_API_KEY and OPENAI_API_KEY in your .env file or system environment.")
         sys.exit(1)
         
    asyncio.run(main())