#Generated python
#!/usr/bin/env python3
"""
Gemini-Enhanced PowerPoint Generator
This script uses Gemini's capabilities to generate rich content and create a PowerPoint presentation.
"""
import asyncio
import sys
import os
from pathlib import Path
import logging
import base64
import uuid
import mimetypes
from io import BytesIO
import tempfile
import shutil
import datetime

# Ensure the autobyteus package is discoverable
SCRIPT_DIR = Path(__file__).resolve().parent
PACKAGE_ROOT = SCRIPT_DIR.parent
sys.path.insert(0, str(PACKAGE_ROOT))

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(name)s - %(message)s')
logger = logging.getLogger("gemini_pptx_generator")

# Import necessary modules
from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.llm.utils.llm_config import LLMConfig

# Create temporary directory for generated content
TEMP_DIR = tempfile.mkdtemp(prefix="gemini_pptx_")

def save_content_to_file(content_bytes, file_extension):
    """
    Saves content bytes to a file
    
    Args:
        content_bytes (bytes): The binary content to save
        file_extension (str): The file extension (e.g., 'png', 'mp4')
        
    Returns:
        str: Path to the saved file
    """
    filename = f"{uuid.uuid4()}.{file_extension}"
    filepath = os.path.join(TEMP_DIR, filename)
    
    with open(filepath, "wb") as f:
        f.write(content_bytes)
        
    logger.info(f"Saved content to {filepath}")
    return filepath

async def generate_title_slide_content(gemini, title):
    """Generate content for the title slide"""
    prompt = f"Create a professional, visually appealing title slide for a presentation titled '{title}'. Include a subtitle that briefly explains what the presentation is about."
    
    try:
        response = await gemini.generate_content(prompt)
        subtitle = response.text.strip().split('\n')[0][:100]  # Take first line, limit to 100 chars
        return subtitle
    except Exception as e:
        logger.error(f"Error generating title slide content: {e}")
        return "A presentation created with Gemini AI"

async def generate_image_content(gemini, prompt="A futuristic city with advanced technology, flying vehicles, and sustainable architecture"):
    """Generate an image using Gemini"""
    try:
        image_bytes_list = await gemini.generate_image(prompt, number_of_images=1)
        if image_bytes_list:
            image_path = save_content_to_file(image_bytes_list[0], "png")
            return image_path, prompt
        else:
            raise ValueError("Failed to generate image")
    except Exception as e:
        logger.error(f"Error generating image: {e}")
        return None, prompt

async def generate_diagram_content(gemini, prompt="A simple software architecture with frontend, backend, and database layers", diagram_type="architecture"):
    """Generate a diagram using Gemini"""
    try:
        diagram_bytes = await gemini.generate_diagram(prompt, diagram_type)
        diagram_path = save_content_to_file(diagram_bytes, "png")
        return diagram_path, prompt
    except Exception as e:
        logger.error(f"Error generating diagram: {e}")
        return None, prompt

async def generate_chart_content(gemini, prompt="Monthly sales data for Q1 and Q2 with clear labels"):
    """Generate a chart using Gemini"""
    try:
        chart_bytes = await gemini.generate_diagram(prompt, diagram_type="chart")
        chart_path = save_content_to_file(chart_bytes, "png")
        return chart_path, prompt
    except Exception as e:
        logger.error(f"Error generating chart: {e}")
        return None, prompt

async def generate_table_content(gemini, prompt="Create a table showing quarterly sales data for different product categories"):
    """Generate table data using Gemini"""
    try:
        response = await gemini.generate_content(prompt + ". Format the data as a table with 5-7 rows and 3-5 columns. Include headers.")
        table_text = response.text
        
        # Parse the table text into rows and columns
        rows = []
        for line in table_text.strip().split('\n'):
            if line.strip() and '|' in line:
                # Split by | and strip each cell
                cells = [cell.strip() for cell in line.split('|')]
                # Remove empty cells that might occur at the beginning or end due to leading/trailing |
                cells = [cell for cell in cells if cell]
                if cells:
                    rows.append(cells)
        
        return rows, prompt
    except Exception as e:
        logger.error(f"Error generating table content: {e}")
        # Return a simple default table
        return [
            ["Category", "Q1", "Q2", "Q3", "Q4"],
            ["Product A", "$10,000", "$12,500", "$15,000", "$17,500"],
            ["Product B", "$8,500", "$9,000", "$11,000", "$12,000"],
            ["Product C", "$5,000", "$5,500", "$6,000", "$7,500"]
        ], prompt

async def create_powerpoint_presentation(title, output_file, pro_gemini, flash_gemini):
    """
    Creates a PowerPoint presentation with Gemini-generated content
    
    Args:
        title (str): The title for the presentation
        output_file (str): Path where to save the PowerPoint file
        pro_gemini: The Gemini Pro LLM instance for text, diagrams, and charts.
        flash_gemini: The Gemini Flash LLM instance for images.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    logger.info(f"Creating PowerPoint presentation: {title}")
    
    # Create a new presentation
    prs = Presentation()
    
    # 1. Create title slide
    logger.info("Creating title slide...")
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle = await generate_title_slide_content(pro_gemini, title)
    subtitle_shape.text = subtitle
    
    # 2. Create agenda slide
    logger.info("Creating agenda slide...")
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "Agenda"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "This presentation includes:"
    
    p = tf.add_paragraph()
    p.text = "AI-Generated Visuals"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Architecture Diagram"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Visualization"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Tables"
    p.level = 1
    
    # 3. Create image slide
    logger.info("Creating AI-generated image slide...")
    image_path, image_prompt = await generate_image_content(flash_gemini)
    if image_path:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = "AI-Generated Image"
        title_para.font.bold = True
        title_para.font.size = Pt(32)
        
        # Add image
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))
        
        # Add caption
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.add_paragraph()
        caption_para.text = image_prompt
        caption_para.font.italic = True
        caption_para.alignment = PP_ALIGN.CENTER
    
    # 4. Create diagram slide
    logger.info("Creating architecture diagram slide...")
    diagram_path, diagram_prompt = await generate_diagram_content(pro_gemini)
    if diagram_path:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = "Architecture Diagram"
        title_para.font.bold = True
        title_para.font.size = Pt(32)
        
        # Add diagram
        slide.shapes.add_picture(diagram_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))
        
        # Add caption
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.add_paragraph()
        caption_para.text = diagram_prompt
        caption_para.font.italic = True
        caption_para.alignment = PP_ALIGN.CENTER
    
    # 5. Create chart slide
    logger.info("Creating data visualization slide...")
    chart_path, chart_prompt = await generate_chart_content(pro_gemini)
    if chart_path:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = "Data Visualization"
        title_para.font.bold = True
        title_para.font.size = Pt(32)
        
        # Add chart
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))
        
        # Add caption
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.add_paragraph()
        caption_para.text = chart_prompt
        caption_para.font.italic = True
        caption_para.alignment = PP_ALIGN.CENTER
    
    # 6. Create table slide
    logger.info("Creating data table slide...")
    table_data, table_prompt = await generate_table_content(pro_gemini)
    if table_data:
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = "Data Table"
        
        # Determine table dimensions
        rows = len(table_data)
        cols = max(len(row) for row in table_data)
        
        # Add table
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Populate the table
        for i, row in enumerate(table_data):
            for j, cell_text in enumerate(row):
                if j < cols:  # Ensure we don't exceed the number of columns
                    cell = table.cell(i, j)
                    cell.text = cell_text
                    
                    # Format header row
                    if i == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
    
    # 7. Create a closing slide
    logger.info("Creating closing slide...")
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "Thank You!"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "This presentation was created with:"
    
    p = tf.add_paragraph()
    p.text = "Gemini AI for content generation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Python and python-pptx for presentation assembly"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Generated on {datetime.datetime.now().strftime('%Y-%m-%d')}"
    p.level = 1
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(os.path.abspath(output_file))
    os.makedirs(output_dir, exist_ok=True)
    
    # Save the presentation
    prs.save(output_file)
    logger.info(f"PowerPoint presentation saved as: {output_file}")
    return True

async def main():
    """Main function to run the example"""
    import argparse
    
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="Create a PowerPoint presentation with Gemini-generated content")
    parser.add_argument("title", nargs="?", default="Gemini-Enhanced Presentation", help="Title of the presentation")
    parser.add_argument("--output-dir", default="./presentations", help="Directory to save the PowerPoint file (default: ./presentations)")
    parser.add_argument("--save-content", action="store_true", help="Save all generated content (images, diagrams) to the output directory")
    args = parser.parse_args()
    
    pro_gemini = None
    flash_gemini = None
    try:
        # This script acts as an automated agent, using specialized LLM models for different tasks.
        logger.info("Initializing Gemini LLMs for the presentation agent...")
        
        # Use Gemini 2.5 Pro for complex tasks like text, diagram, and chart generation
        logger.info("Initializing Gemini 2.5 Pro for text and reasoning...")
        pro_gemini = LLMFactory.create_llm(
            model_identifier="gemini-2.5-pro",
            llm_config=LLMConfig(temperature=0.7)
        )
        
        # Use the faster Gemini 2.0 Flash for image generation
        logger.info("Initializing Gemini 2.0 Flash for image generation...")
        flash_gemini = LLMFactory.create_llm(
            model_identifier="gemini-2.0-flash",
            llm_config=LLMConfig(temperature=0.7)
        )
        
        # Create output directory if it doesn't exist
        output_dir = os.path.abspath(args.output_dir)
        os.makedirs(output_dir, exist_ok=True)
        
        # Set output file path
        safe_title = args.title.replace(' ', '_').replace('/', '_').replace('\\', '_')
        output_file = os.path.join(output_dir, f"{safe_title}.pptx")
        
        # Create the PowerPoint presentation
        success = await create_powerpoint_presentation(args.title, output_file, pro_gemini, flash_gemini)
        
        if success:
            print(f"\nPresentation saved as: {output_file}")
            
            # Copy generated content to output directory if requested
            if args.save_content:
                content_dir = os.path.join(output_dir, "content")
                os.makedirs(content_dir, exist_ok=True)
                
                # Copy all generated content from temp directory to output directory
                for file in os.listdir(TEMP_DIR):
                    if file.endswith(('.png', '.jpg', '.jpeg', '.mp4', '.gif')):
                        src_file = os.path.join(TEMP_DIR, file)
                        dst_file = os.path.join(content_dir, file)
                        shutil.copy2(src_file, dst_file)
                        logger.info(f"Copied generated content: {dst_file}")
                
                print(f"Generated content saved to: {content_dir}")
        else:
            print("\nFailed to create presentation.")
    except Exception as e:
        logger.error(f"Error creating presentation: {e}", exc_info=True)
    finally:
        # Clean up
        if pro_gemini:
            await pro_gemini.cleanup()
        if flash_gemini:
            await flash_gemini.cleanup()

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("\nScript interrupted by user.")
    except Exception as e:
        logger.error(f"Unhandled error: {e}", exc_info=True)
        sys.exit(1)