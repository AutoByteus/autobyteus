import os
import sys
import asyncio
import logging
import json
import uuid
import aiohttp
import pptx
from pptx.util import Inches, Pt
from typing import List, Dict, Any

# Ensure the project root is in the Python path
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.tools.functional_tool import tool
from autobyteus.agent.context.agent_config import AgentConfig
from autobyteus.agent.group.agent_group import AgentGroup
from autobyteus.utils.file_utils import get_default_download_folder
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# --- Directory Setup ---
IMAGE_DIR = os.path.join(get_default_download_folder(), "generated_ppt_images")
PPT_DIR = os.path.join(get_default_download_folder(), "generated_ppts")
os.makedirs(IMAGE_DIR, exist_ok=True)
os.makedirs(PPT_DIR, exist_ok=True)

# --- Tool Definitions ---

@tool(name="ImageGenerator")
async def generate_image(prompt: str, image_name: str = "generated_image") -> str:
    """
    Generates an image using DALL-E 3 based on a descriptive prompt,
    downloads it, and saves it to a local directory.

    Args:
        prompt: A detailed English description of the image to be generated.
        image_name: A base name for the file (e.g., 'solar_panels'). A unique ID will be appended.

    Returns:
        The local file path of the saved image.
    """
    logger.info(f"ImageGenerator tool called with prompt: '{prompt[:50]}...'")
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        error_msg = "OPENAI_API_KEY environment variable not set. Cannot generate image."
        logger.error(error_msg)
        raise ValueError(error_msg)

    try:
        # Use official openai library if available
        from openai import AsyncOpenAI
        client = AsyncOpenAI(api_key=api_key)
        response = await client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        image_url = response.data[0].url
    except ImportError:
        # Fallback to manual aiohttp request if openai library is not installed
        logger.warning("openai library not found, using aiohttp for DALL-E request.")
        async with aiohttp.ClientSession() as session:
            async with session.post(
                "https://api.openai.com/v1/images/generations",
                headers={"Authorization": f"Bearer {api_key}"},
                json={"model": "dall-e-3", "prompt": prompt, "n": 1, "size": "1024x1024"}
            ) as resp:
                data = await resp.json()
                if not resp.ok:
                    raise Exception(f"DALL-E API error: {data.get('error', {}).get('message', 'Unknown error')}")
                image_url = data['data'][0]['url']

    logger.info(f"Image generated successfully. URL: {image_url}")

    # Download the generated image
    async with aiohttp.ClientSession() as session:
        async with session.get(image_url) as resp:
            if resp.status != 200:
                error_text = await resp.text()
                raise Exception(f"Failed to download image from {image_url}. Status: {resp.status}. Response: {error_text}")
            image_bytes = await resp.read()

    # Save the image locally
    safe_image_name = "".join(c for c in image_name if c.isalnum() or c in (' ', '_')).rstrip()
    unique_filename = f"{safe_image_name}_{uuid.uuid4().hex[:8]}.png"
    file_path = os.path.join(IMAGE_DIR, unique_filename)
    with open(file_path, "wb") as f:
        f.write(image_bytes)
    logger.info(f"Image saved locally to: {file_path}")
    return file_path

@tool(name="PPTCompiler")
def compile_presentation(slides_data_json: str, presentation_title: str) -> str:
    """
    Compiles a PowerPoint presentation from a JSON string of slide data.

    Args:
        slides_data_json: A JSON string representing a list of slides. Each slide
                          is a dictionary with "title", "content", and "image_path".
        presentation_title: The main title for the presentation.

    Returns:
        The local file path of the saved .pptx file.
    """
    logger.info(f"PPTCompiler tool called for presentation: '{presentation_title}'")
    try:
        slides_data: List[Dict[str, Any]] = json.loads(slides_data_json)
    except json.JSONDecodeError as e:
        logger.error(f"Failed to decode slides_data_json: {e}")
        raise ValueError("Invalid JSON format for slides_data_json.") from e

    prs = pptx.Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders.get(1)
    title.text = presentation_title
    if subtitle:
        subtitle.text = "Generated by Multi-Agent System"

    # Content Slides
    for slide_data in slides_data:
        layout = prs.slide_layouts[5]  # Title and Content layout
        slide = prs.slides.add_slide(layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Untitled Slide")

        # Content
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(5.0)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        content = slide_data.get("content", "")
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(14)

        # Image
        image_path = slide_data.get("image_path")
        if image_path and os.path.exists(image_path):
            try:
                img_left = Inches(5.75)
                img_top = Inches(1.8)
                img_height = Inches(3.5)
                slide.shapes.add_picture(image_path, img_left, img_top, height=img_height)
                logger.info(f"Added image {image_path} to slide '{title_shape.text}'")
            except Exception as e:
                logger.warning(f"Could not add image {image_path} to slide: {e}")
        elif image_path:
            logger.warning(f"Image path specified but not found: {image_path}")


    # Save presentation
    safe_title = "".join(c for c in presentation_title if c.isalnum() or c in (' ', '_')).rstrip()
    filename = f"{safe_title}_{uuid.uuid4().hex[:8]}.pptx"
    file_path = os.path.join(PPT_DIR, filename)
    prs.save(file_path)
    logger.info(f"Presentation saved successfully to: {file_path}")
    return file_path

# --- System Prompts for Agents ---

outline_agent_system_prompt = """You are a master of creating presentation outlines.
Based on the user's request for a presentation topic, your task is to generate a slide-by-slide outline.
The output must be a single JSON object with a key "slides", which is a list of slide objects.
Each slide object in the list must have a "title" key.
Example for a 2-slide topic: {"slides": [{"title": "Slide 1 Title"}, {"title": "Slide 2 Title"}]}
You must only output the JSON object and nothing else. No preamble, no explanation."""

content_agent_system_prompt = """You are an expert content creator for presentations.
You will receive a single slide title as input.
Your tasks are:
1.  Write engaging and informative content for this slide title. The content should be a concise paragraph of 2-4 sentences.
2.  Based on the slide's title and content, create a detailed and visually descriptive prompt for an image generation model (DALL-E 3). The prompt should be in English and aim for a professional, photorealistic, or concept art style suitable for a presentation.
3.  Call the "ImageGenerator" tool with the image prompt you just created. Give the image a descriptive name based on the slide title.
4.  Finally, output a single JSON object containing three keys: "title" (the original input title), "content" (the text you wrote), and "image_path" (the local file path returned by the ImageGenerator tool).

You must only output the final JSON object after the tool has been called. Do not add any other text or explanation.
"""

compiler_agent_system_prompt = """You are the final PowerPoint assembler.
You will receive a JSON string containing the presentation title and a list of all slide data.
Each slide in the list has a 'title', 'content', and 'image_path'.
Your ONLY task is to call the "PPTCompiler" tool with the provided data.
Do not modify, add, or analyze the content. Just call the tool."""

coordinator_agent_system_prompt = """You are a project manager responsible for creating a PowerPoint presentation. You will orchestrate a team of agents to complete this task.

Your team consists of:
- OutlineAgent (Role: 'Presentation Structurer'): Generates the presentation outline.
- ContentAgent (Role: 'Slide Content and Image Prompt Generator'): Creates content and an image for a single slide.
- CompilerAgent (Role: 'PowerPoint Assembler'): Compiles the final .pptx file.

Your workflow is as follows:
1.  You will receive the main topic for the presentation from the user.
2.  **Delegate to OutlineAgent**: Use the 'SendMessageTo' tool to send the topic to the 'Presentation Structurer' agent. Wait for the JSON outline as a response.
3.  **Process Each Slide**:
    - Once you receive the outline (which will be a string containing JSON), parse it. Then, iterate through each slide title in the JSON response.
    - For each slide title, use the 'SendMessageTo' tool to delegate the task to the 'Slide Content and Image Prompt Generator' agent. Send only the slide title as the content of the message.
    - The ContentAgent will return a JSON object (as a string) for that single slide, including the 'image_path'.
4.  **Aggregate Results**: Collect the JSON objects for all the processed slides into a single list.
5.  **Delegate to CompilerAgent**:
    - Once all slides are processed, create a final JSON object with two keys: "presentation_title" and "slides_data" (which is the list you just aggregated).
    - Use the 'SendMessageTo' tool to send this complete JSON object as a string to the 'PowerPoint Assembler' agent.
6.  **Finalize**: The CompilerAgent will return the file path of the final .pptx file. Your final response to the user should be ONLY this file path and nothing else.

IMPORTANT: You must use the 'SendMessageTo' tool for all communication between agents. The `recipient_role_name` parameter is how you select which agent to talk to.
"""

# --- Agent Configurations ---

outline_agent_config = AgentConfig(
    name="OutlineAgent",
    role="Presentation Structurer",
    description="Generates a slide-by-slide outline from a user request.",
    llm_instance=LLMFactory.create_llm("GPT_4o_API"),
    system_prompt=outline_agent_system_prompt,
    tools=[],
    use_xml_tool_format=False,
)

content_agent_config = AgentConfig(
    name="ContentAgent",
    role="Slide Content and Image Prompt Generator",
    description="Takes a single slide topic, writes content, generates an image prompt, and creates the image.",
    llm_instance=LLMFactory.create_llm("GPT_4o_API"),
    system_prompt=content_agent_system_prompt,
    tools=[generate_image],
    use_xml_tool_format=False,
)

compiler_agent_config = AgentConfig(
    name="CompilerAgent",
    role="PowerPoint Assembler",
    description="Gathers all slide data and compiles the final .pptx file.",
    llm_instance=LLMFactory.create_llm("GPT_4o_API"),
    system_prompt=compiler_agent_system_prompt,
    tools=[compile_presentation],
    use_xml_tool_format=False,
)

coordinator_agent_config = AgentConfig(
    name="CoordinatorAgent",
    role="Project Manager",
    description="Orchestrates the entire presentation generation workflow.",
    llm_instance=LLMFactory.create_llm("GPT_4o_API"),
    system_prompt=coordinator_agent_system_prompt,
    tools=[],  # The SendMessageTo tool is added automatically by the AgentGroup
    use_xml_tool_format=False,
)


# --- Main Workflow Execution ---

async def run_ppt_generation_workflow(topic: str):
    """Initializes and runs the multi-agent PowerPoint generation workflow."""
    logger.info("--- PowerPoint Generation Workflow Initializing ---")
    
    agent_group = AgentGroup(
        agent_configs=[
            outline_agent_config,
            content_agent_config,
            compiler_agent_config,
            coordinator_agent_config, 
        ],
        coordinator_config_name="CoordinatorAgent"
    )

    try:
        logger.info("Starting the agent group...")
        await agent_group.start()
        
        logger.info(f"Sending initial topic to coordinator: '{topic}'")
        final_result = await agent_group.process_task_for_coordinator(
            initial_input_content=topic
        )
        
        logger.info("--- Workflow Completed ---")
        logger.info(f"Final result from coordinator: {final_result}")
        
        print("\n" + "="*50)
        print("PowerPoint Generation Complete!")
        print(f"Presentation saved to: {final_result}")
        print("="*50 + "\n")
        
    except Exception as e:
        logger.error(f"An error occurred during the workflow: {e}", exc_info=True)
    finally:
        logger.info("Stopping the agent group...")
        await agent_group.stop()
        logger.info("Workflow has been shut down.")

if __name__ == "__main__":
    if not os.getenv("OPENAI_API_KEY"):
        print("\nERROR: OPENAI_API_KEY environment variable is not set.")
        print("Please set it in a .env file or as an environment variable.")
        print("e.g., export OPENAI_API_KEY='sk-...'")
        sys.exit(1)

    if len(sys.argv) > 1:
        user_topic = " ".join(sys.argv[1:])
    else:
        user_topic = "A 4-slide presentation on the future of renewable energy, covering solar, wind, and geothermal."
        print(f"No topic provided. Using default topic: '{user_topic}'")
    
    asyncio.run(run_ppt_generation_workflow(topic=user_topic))