import os
import logging
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches
from autobyteus.tools import tool
from autobyteus.tools.tool_category import ToolCategory

logger = logging.getLogger(__name__)

@tool(name="PowerPointCreator", category=ToolCategory.FILE_SYSTEM)
async def create_powerpoint(
    title: str, slides_data: List[Dict[str, Any]], output_path: str
) -> str:
    """
    Creates a PowerPoint presentation.

    Args:
        title: The title of the presentation.
        slides_data: A list of dictionaries, where each dictionary represents a slide.
        output_path: The path to save the presentation to.

    Returns:
        The path to the saved presentation.
    """
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]

        content_shape = slide.placeholders[1]
        content_shape.text = slide_data["content"]

        if "image_path" in slide_data:
            img_path = slide_data["image_path"]
            left = Inches(1)
            top = Inches(3)
            slide.shapes.add_picture(img_path, left, top, height=Inches(4))

    prs.save(output_path)
    return output_path