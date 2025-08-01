import os
import json
import logging
import re
from typing import Dict, List, Optional, Any
from datetime import datetime

from autobyteus.tools import tool
from autobyteus.tools.tool_category import ToolCategory
from pptx import Presentation
from pptx.util import Pt, Inches

if TYPE_CHECKING:
    from autobyteus.agent.context import AgentContext

logger = logging.getLogger(__name__)

OUTPUT_DIR = "generated_presentations"
os.makedirs(OUTPUT_DIR, exist_ok=True)

@tool(name="PPTSaverTool", category=ToolCategory.FILE_SYSTEM)
async def ppt_saver_tool(context: 'AgentContext', presentation_json: str) -> str:
    """
    Takes a JSON string representing a presentation, generates a .pptx file,
    and returns its public-facing URL.
    """
    logger.info(f"Tool 'PPTSaverTool' called by agent '{context.agent_id}'.")
    try:
        data = json.loads(presentation_json)
        if isinstance(data, list):
            # If the LLM returns a list of slide objects
            data = {"sections": data, "title": "Presentation"}
    except json.JSONDecodeError as e:
        error_msg = f"Invalid JSON format provided to PPTSaverTool: {e}"
        logger.error(error_msg)
        return f"Error: {error_msg}"

    try:
        prs = Presentation()
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = data.get("title", "Presentation")
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

        # Content Slides
        for section in data.get("sections", []):
            content_slide_layout = prs.slide_layouts[5] # Blank layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Title
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            title_text = "Slide Title"
            content_data = section.get("content", [])
            if content_data and content_data[0].get("type") == "h1":
                title_text = content_data[0].get("children", [{}])[0].get("text", "Slide Title")
            p.text = title_text
            p.font.bold = True
            p.font.size = Pt(32)

            # Body
            body_text = ""
            for item in content_data:
                 if item.get("type") == "bullets":
                     for bullet in item.get("children", []):
                         h3 = bullet.get("children", [{}])[0].get("children", [{}])[0].get("text", "")
                         p_text = bullet.get("children", [{}])[1].get("children", [{}])[0].get("text", "")
                         body_text += f"• {h3}: {p_text}\n"
            
            txBox_body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            tf_body = txBox_body.text_frame
            tf_body.text = body_text
            tf_body.word_wrap = True

        sanitized_title = re.sub(r'[\\/*?:"<>|]', "", data.get("title", "presentation"))
        filename = f"{sanitized_title}_{int(datetime.now().timestamp())}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)

        public_url = f'{os.environ.get("PUBLIC_HOST_URL", "http://localhost:8888")}/downloads/{filename}'
        logger.info(f"Presentation saved to '{filepath}'. Public URL: {public_url}")
        return public_url
    except Exception as e:
        logger.error(f"Failed to generate PPTX file: {e}", exc_info=True)
        return f"Error: Failed to generate presentation. Reason: {str(e)}"