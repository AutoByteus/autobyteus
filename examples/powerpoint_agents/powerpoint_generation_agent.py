#!/usr/bin/env python3
import asyncio
import logging
import argparse
from pathlib import Path
import sys
import os
import json
import tempfile
import shutil
from typing import Dict, Any, List
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).resolve().parent
PACKAGE_ROOT = SCRIPT_DIR.parent
if str(PACKAGE_ROOT) not in sys.path:
    sys.path.insert(0, str(PACKAGE_ROOT))

try:
    from dotenv import load_dotenv
    env_file_path = PACKAGE_ROOT / ".env"
    if env_file_path.exists():
        load_dotenv(env_file_path)
except ImportError:
    pass

from autobyteus.agent.context.agent_config import AgentConfig
from autobyteus.agent.factory.agent_factory import AgentFactory
from autobyteus.agent.message.agent_input_user_message import AgentInputUserMessage
from autobyteus.llm.models import LLMModel
from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.llm.utils.llm_config import LLMConfig
from autobyteus.tools.base_tool import BaseTool
from autobyteus.tools.functional_tool import FunctionalTool, tool
from autobyteus.tools.registry.tool_registry import ToolRegistry

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

TEMP_DIR = tempfile.mkdtemp(prefix="pptx_agent_")

# Initialize Gemini models globally
gemini_pro = None
gemini_flash = None

async def init_gemini_models():
    global gemini_pro, gemini_flash
    if not gemini_pro:
        gemini_pro = LLMFactory.create_llm(
            model_identifier="gemini-2.5-pro",
            llm_config=LLMConfig(temperature=0.7)
        )
    if not gemini_flash:
        gemini_flash = LLMFactory.create_llm(
            model_identifier="gemini-2.0-flash",
            llm_config=LLMConfig(temperature=0.7)
        )

async def generate_slide_content(user_input: str, num_slides: int = 5) -> Dict[str, Any]:
    """Generate structured slide content using Gemini Pro"""
    await init_gemini_models()
    
    prompt = f"""Create a PowerPoint presentation for: "{user_input}"
    
Generate exactly {num_slides} slides in this JSON format:
{{
    "title": "Main Presentation Title",
    "slides": [
        {{
            "number": 1,
            "type": "title",
            "title": "Title text",
            "subtitle": "Subtitle text",
            "bullets": [],
            "notes": "Speaker notes",
            "image_prompt": "Detailed image generation prompt"
        }},
        {{
            "number": 2,
            "type": "content",
            "title": "Slide title",
            "bullets": ["Point 1", "Point 2", "Point 3"],
            "notes": "Speaker notes",
            "image_prompt": "Detailed image generation prompt"
        }}
    ]
}}

Make it professional and informative. Return ONLY valid JSON."""
    
    response = await gemini_pro.generate_content(prompt)
    
    try:
        content = json.loads(response.text.strip())
        logger.info(f"Generated {len(content['slides'])} slides")
        return content
    except:
        return {
            "title": user_input,
            "slides": [{
                "number": 1,
                "type": "title",
                "title": user_input,
                "subtitle": "AI-Generated Presentation",
                "bullets": [],
                "notes": "",
                "image_prompt": f"Professional image for {user_input}"
            }]
        }

async def generate_slide_image(prompt: str, slide_num: int) -> str:
    """Generate image using Gemini Flash"""
    await init_gemini_models()
    
    try:
        response = await gemini_flash.generate_content(
            prompt,
            config={'response_modalities': ['IMAGE']}
        )
        
        for part in response.candidates[0].content.parts:
            if hasattr(part, 'inline_data') and part.inline_data:
                image = Image.open(BytesIO(part.inline_data.data))
                filepath = os.path.join(TEMP_DIR, f"slide_{slide_num}.png")
                image.save(filepath)
                logger.info(f"Generated image for slide {slide_num}")
                return filepath
    except Exception as e:
        logger.error(f"Image generation failed: {e}")
    
    # Create placeholder
    image = Image.new('RGB', (800, 600), 'lightgray')
    filepath = os.path.join(TEMP_DIR, f"slide_{slide_num}_placeholder.png")
    image.save(filepath)
    return filepath

async def create_powerpoint(content: Dict[str, Any], images: Dict[int, str], output_path: str) -> str:
    """Create PowerPoint file"""
    prs = Presentation()
    
    for slide_data in content['slides']:
        num = slide_data['number']
        
        if slide_data['type'] == 'title' and num == 1:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = slide_data['title']
            if slide_data.get('subtitle'):
                slide.placeholders[1].text = slide_data['subtitle']
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
            title_box.text_frame.paragraphs[0].text = slide_data['title']
            title_box.text_frame.paragraphs[0].font.size = Pt(28)
            title_box.text_frame.paragraphs[0].font.bold = True
            
            # Content
            if slide_data.get('bullets'):
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
                tf = content_box.text_frame
                for i, bullet in enumerate(slide_data['bullets']):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(18)
            
            # Image
            if num in images and os.path.exists(images[num]):
                slide.shapes.add_picture(
                    images[num],
                    Inches(5.5), Inches(1.5),
                    width=Inches(4), height=Inches(4.5)
                )
        
        # Notes
        if slide_data.get('notes'):
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    logger.info(f"Saved: {output_path}")
    return output_path

# Create tools using the @tool decorator
@tool(name="generate_slides", description="Generate slide content from user input")
async def generate_slides_tool(user_input: str, num_slides: int = 5) -> str:
    result = await generate_slide_content(user_input, num_slides)
    import json
    return json.dumps(result)

@tool(name="generate_image", description="Generate image for a slide")
async def generate_image_tool(prompt: str, slide_num: int) -> str:
    return await generate_slide_image(prompt, slide_num)

@tool(name="create_pptx", description="Create PowerPoint from content and images")
async def create_pptx_tool(content: str, images: str, output_path: str) -> str:
    # Parse JSON strings back to objects
    import json
    content_dict = json.loads(content) if isinstance(content, str) else content
    images_raw = json.loads(images) if isinstance(images, str) else images
    # Convert string keys to integers for images dict
    images_dict = {int(k): v for k, v in images_raw.items()}
    return await create_powerpoint(content_dict, images_dict, output_path)

def create_tools():
    return [generate_slides_tool, generate_image_tool, create_pptx_tool]

SYSTEM_PROMPT = """You are a PowerPoint Generation Agent. You MUST use the provided tools to actually create presentations.

CRITICAL: You must invoke tools using this exact XML format. Do NOT just describe what you would do.

Example:
<tool name="generate_slides">
    <arguments>
        <arg name="user_input">AI in healthcare</arg>
        <arg name="num_slides">3</arg>
    </arguments>
</tool>

WORKFLOW:
1. Start by calling generate_slides with the user's topic and number of slides
2. When you get the JSON response, identify slides that have image_prompt
3. For each slide needing an image, call generate_image with the image_prompt
4. Finally, call create_pptx with the content JSON and images dictionary

For create_pptx:
- content: the exact JSON string from generate_slides
- images: JSON like {"1": "/path/to/image.png", "2": "/path/to/image2.png"}  
- output_path: the specified file path

You MUST call the tools, not just plan what you would do."""

async def main():
    parser = argparse.ArgumentParser(description="PowerPoint Generation Agent")
    parser.add_argument("input", help="Title or paragraph for presentation")
    parser.add_argument("--output", default="./presentation.pptx", help="Output path")
    parser.add_argument("--slides", type=int, default=3, help="Number of slides")
    parser.add_argument("--agent-llm", default="gemini-2.5-pro", help="LLM for agent orchestration")
    
    args = parser.parse_args()
    
    try:
        # Initialize agent
        logger.info("Initializing PowerPoint Agent...")
        
        agent_llm = LLMFactory.create_llm(
            model_identifier=args.agent_llm,
            llm_config=LLMConfig(temperature=0.3)
        )
        
        tools = create_tools()
        
        config = AgentConfig(
            name="PowerPointAgent",
            role="Presentation Creator",
            description="Creates PowerPoint presentations using Gemini",
            llm_instance=agent_llm,
            system_prompt=SYSTEM_PROMPT,
            tools=tools,
            auto_execute_tools=True
        )
        
        agent = AgentFactory().create_agent(config=config)
        
        # Send request
        message_content = f"""Create a presentation about: {args.input}
Number of slides: {args.slides}
Output path: {args.output}

Generate all content and images, then create the PowerPoint file."""
        
        user_message = AgentInputUserMessage(content=message_content)
        
        # Start the agent if not already running
        agent.start()
        await asyncio.sleep(0.5)  # Give it time to start
        
        await agent.post_user_message(user_message)
        
        # Wait for agent to complete all work
        from autobyteus.agent.utils.wait_for_idle import wait_for_agent_to_be_idle
        await wait_for_agent_to_be_idle(agent, timeout=120.0)  # Longer timeout for tool execution
        
        logger.info(f"✅ Presentation created: {args.output}")
        
        # Save images separately
        output_dir = Path(args.output).parent
        images_dir = output_dir / "images"
        images_dir.mkdir(exist_ok=True)
        
        for file in os.listdir(TEMP_DIR):
            if file.endswith('.png'):
                shutil.copy2(
                    os.path.join(TEMP_DIR, file),
                    images_dir / file
                )
        
        logger.info(f"📁 Images saved to: {images_dir}")
        
    except Exception as e:
        logger.error(f"Error: {e}", exc_info=True)
    finally:
        # Cleanup
        if gemini_pro:
            await gemini_pro.cleanup()
        if gemini_flash:
            await gemini_flash.cleanup()
        shutil.rmtree(TEMP_DIR, ignore_errors=True)

if __name__ == "__main__":
    asyncio.run(main()) 