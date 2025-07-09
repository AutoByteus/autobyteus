#!/usr/bin/env python
# -*- coding: utf-8 -*-

import asyncio
import logging
import argparse
from pathlib import Path
import sys
import os
import json
import tempfile
import shutil
import httpx
import openai
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from typing import Optional, Any

# --- Core System Setup ---
# Add the project root to the Python path to ensure all modules are found.
# This is necessary because this script is designed to be a self-contained example.
SCRIPT_DIR = Path(__file__).resolve().parent
# Assuming the script is in 'src/examples', the project root is two levels up.
# Adjust if the script's location is different relative to the 'autobyteus' and 'llm' packages.
PROJECT_ROOT = SCRIPT_DIR.parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

# --- Environment and Logging Configuration ---
# Load environment variables from a .env file if it exists.
# This is where API keys (OPENAI_API_KEY) should be stored.
try:
    from dotenv import load_dotenv
    # Assuming .env file is in the project root
    env_file_path = PROJECT_ROOT / ".env"
    if env_file_path.exists():
        print(f"Loading environment variables from: {env_file_path}")
        load_dotenv(dotenv_path=env_file_path)
    else:
        print(f"Info: .env file not found at {env_file_path}. Relying on system environment variables.")
except ImportError:
    print("Warning: python-dotenv not installed. Relying on system environment variables.")
    pass

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("PowerPointAgentApp")

# --- Autobyteus Framework Imports ---
# Import necessary components from the provided autobyteus framework files.
from autobyteus.agent.context.agent_config import AgentConfig
from autobyteus.agent.factory.agent_factory import AgentFactory
from autobyteus.agent.message.agent_input_user_message import AgentInputUserMessage
from autobyteus.agent.utils.wait_for_idle import wait_for_agent_to_be_idle
from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.llm.utils.llm_config import LLMConfig
from autobyteus.tools.functional_tool import tool
from autobyteus.agent.runtime.agent_thread_pool_manager import AgentThreadPoolManager

# --- Global State and Setup ---
# Create a temporary directory for storing generated images.
TEMP_DIR = Path(tempfile.mkdtemp(prefix="pptx_agent_"))
logger.info(f"Temporary directory for images created at: {TEMP_DIR}")

# --- Tool Definitions ---
# The agent will use these tools to perform its tasks. The `@tool` decorator
# automatically registers them with the framework and generates a schema for the LLM.

@tool(name="generate_slide_content")
async def generate_slide_content(user_input: str, num_slides: int = 5) -> str:
    """
    Generates structured JSON content for a presentation based on a user's topic.

    Args:
        user_input: The central topic or title for the presentation.
        num_slides: The desired number of content slides, in addition to the title slide.

    Returns:
        A JSON string containing the presentation title and a list of slide objects.
        Each slide object has a number, type, title, subtitle, bullets, notes, and an image_prompt.
    """
    logger.info(f"Generating slide content for topic: '{user_input}' ({num_slides} slides)")
    llm = LLMFactory.create_llm("gpt-4o", llm_config=LLMConfig(temperature=0.7))
    prompt = f"""
    Based on the user input "{user_input}", generate the content for a {num_slides}-slide presentation.
    The output must be a single, valid JSON object. Do not include any text or formatting outside of the JSON.
    The JSON structure should be:
    {{
      "presentation_title": "A concise, catchy title for the presentation",
      "slides": [
        {{
          "number": 1,
          "type": "title",
          "title": "The main title of the presentation",
          "subtitle": "A brief, engaging subtitle",
          "notes": "Speaker notes for the title slide.",
          "image_prompt": "A simple, professional DALL-E 3 style prompt for a background image for the title slide."
        }},
        // ... list of {num_slides} more slides of type 'content' ...
        {{
          "number": 2,
          "type": "content",
          "title": "Title for Slide 2",
          "bullets": [
            "A key point or bullet for this slide.",
            "Another key point.",
            "A final key point, keep it concise."
          ],
          "notes": "Speaker notes explaining the content of slide 2 in more detail.",
          "image_prompt": "A simple, professional DALL-E 3 style prompt for an image that visually represents the content of this slide."
        }}
      ]
    }}
    Ensure every slide has a unique `image_prompt`. The prompts should be creative and relevant.
    """
    try:
        response = await llm.send_user_message(prompt)
        # Clean the response to extract only the JSON part
        json_text = response.content.strip().replace("```json", "").replace("```", "").strip()
        # Validate and return the JSON
        json.loads(json_text) # This will raise an error if the JSON is invalid
        logger.info(f"Successfully generated JSON content for {num_slides + 1} slides.")
        return json_text
    except Exception as e:
        logger.error(f"Failed to generate or parse slide content JSON: {e}", exc_info=True)
        # Return a fallback JSON structure on error
        fallback_json = json.dumps({
            "presentation_title": user_input,
            "slides": [{
                "number": 1, "type": "title", "title": user_input,
                "subtitle": "AI-Generated Presentation", "notes": "This is a fallback slide due to an error.",
                "image_prompt": "A generic professional background image for a business presentation."
            }]
        })
        return fallback_json
    finally:
        await llm.cleanup()

@tool(name="generate_slide_image")
async def generate_slide_image(image_prompt: str, slide_number: int) -> str:
    """
    Generates an image using DALL-E 3 based on a text prompt and saves it locally.

    Args:
        image_prompt: A descriptive prompt for the image to be generated.
        slide_number: The number of the slide this image is for, used for naming the file.

    Returns:
        The local file path of the saved image.
    """
    logger.info(f"Generating DALL-E 3 image for slide {slide_number} with prompt: '{image_prompt}'")
    filepath = TEMP_DIR / f"slide_{slide_number}.png"
    
    try:
        # It will automatically use the OPENAI_API_KEY from the environment
        client = openai.AsyncOpenAI()

        logger.info("Requesting image from DALL-E 3...")
        response = await client.images.generate(
            model="dall-e-3",
            prompt=image_prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )

        image_url = response.data[0].url
        if not image_url:
            raise ValueError("DALL-E 3 API did not return an image URL.")
        
        logger.info(f"DALL-E 3 image generated, URL: {image_url}")

        logger.info(f"Downloading image from {image_url}...")
        async with httpx.AsyncClient() as http_client:
            image_response = await http_client.get(image_url, timeout=60.0)
            image_response.raise_for_status()

        with open(filepath, "wb") as f:
            f.write(image_response.content)
            
        logger.info(f"Image downloaded and saved to: {filepath}")
        return str(filepath.resolve())

    except Exception as e:
        logger.error(f"DALL-E 3 image generation failed for slide {slide_number}: {e}", exc_info=True)
        # Create a fallback placeholder image on error
        img = Image.new('RGB', (1024, 1024), (128, 128, 128))
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except IOError:
            font = ImageFont.load_default()
        draw.text((50, 50), f"Slide {slide_number}\nImage Generation Failed", font=font, fill=(255, 255, 255))
        img.save(filepath)
        return str(filepath.resolve())

@tool(name="create_powerpoint_presentation")
async def create_powerpoint_presentation(slide_content_json: str, image_paths_json: str, output_path: str) -> str:
    """
    Creates a PowerPoint .pptx file from structured content and image paths.

    Args:
        slide_content_json: A JSON string containing the presentation title and slide data.
        image_paths_json: A JSON string mapping slide numbers to their local image file paths.
        output_path: The desired file path for the final .pptx presentation.

    Returns:
        A success message with the final path of the created presentation.
    """
    logger.info(f"Creating PowerPoint presentation at: {output_path}")
    try:
        content = json.loads(slide_content_json)
        image_paths = {int(k): v for k, v in json.loads(image_paths_json).items()}
    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON format for content or image paths. Details: {e}"

    prs = Presentation()
    # Set slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for slide_data in content.get('slides', []):
        try:
            if not isinstance(slide_data, dict):
                logger.warning(f"Skipping slide data because it is not a dictionary: {slide_data}")
                continue

            slide_num = slide_data.get('number')
            slide_type = slide_data.get('type')

            if slide_type == 'title':
                slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data.get('title', 'Presentation Title')
                if slide.placeholders[1]:
                    slide.placeholders[1].text = slide_data.get('subtitle', '')
            elif slide_type == 'content':
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)

                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1.2))
                title_frame = title_shape.text_frame
                title_text = slide_data.get('title', '')
                title_frame.text = title_text
                if title_text and title_frame.paragraphs:
                    p = title_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(44)

                has_image = isinstance(slide_num, int) and slide_num in image_paths and Path(image_paths[slide_num]).exists()

                # Add bullets
                bullets = slide_data.get('bullets')
                if bullets and isinstance(bullets, list):
                    text_width = Inches(7.5) if has_image else Inches(15)
                    body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), text_width, Inches(6.5))
                    tf = body_shape.text_frame
                    tf.clear()  # Clear existing paragraph
                    for bullet_text in bullets:
                        p = tf.add_paragraph()
                        p.text = str(bullet_text)
                        p.font.size = Pt(28)
                        p.level = 0

                # Add image
                if has_image:
                    img_path = image_paths[slide_num]
                    slide.shapes.add_picture(img_path, Inches(8), Inches(1.5), width=Inches(7.5))

            # Add speaker notes
            notes_text = slide_data.get('notes')
            if notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text

        except Exception as e:
            logger.error(f"Failed to process slide data: {slide_data}. Error: {e}", exc_info=True)
            continue

    # Ensure output directory exists and save the presentation
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info(f"Presentation saved successfully to {output_path}")
    return f"Presentation created and saved to {output_path}"


# --- Agent System Prompt ---
# This prompt guides the agent on how to use the tools in a sequence.
SYSTEM_PROMPT = """You are a robot that executes a sequence of tool calls. You must not speak or respond in any way other than tool calls.
You will be given a task. You must follow the workflow to complete the task.
You must use the exact XML format for tool calls as shown in the examples provided in the Tool Definition section.

TASK: Create a PowerPoint presentation.
INPUTS: The user will provide the topic, number of slides, and output path.

WORKFLOW:
1. Call `generate_slide_content` with `user_input` and `num_slides` from the user's request.
2. For each slide in the result from step 1, call `generate_slide_image` with the `image_prompt` and `slide_number`.
3. Call `create_powerpoint_presentation` with the JSON from step 1, a JSON map of slide numbers to image paths from step 2, and the `output_path` from the user's request.

The user's request is next. Start with step 1 of the workflow immediately. Do not write any other text.

{{tools}}
"""

# --- Main Execution --
async def get_tool_result(agent: 'Agent', tool_name: str, timeout: int = 180) -> Optional[Any]:
    """Waits for a specific tool call to complete and returns its result."""
    future = asyncio.get_running_loop().create_future()

    def on_event(event: 'AgentEvent'):
        if event.event_type == 'tool_call_result' and event.tool_name == tool_name:
            if not future.done():
                future.set_result(event.result)

    agent.context.event_manager.subscribe(on_event)
    try:
        return await asyncio.wait_for(future, timeout=timeout)
    except asyncio.TimeoutError:
        logger.error(f"Timed out waiting for result from tool '{tool_name}'")
        return None
    finally:
        agent.context.event_manager.unsubscribe(on_event)


async def main():
    """
    Main function to run the PowerPoint generation agent.
    This has been refactored to orchestrate the multi-step process in Python
    for improved reliability.
    """
    parser = argparse.ArgumentParser(description="PowerPoint Generation Agent using GPT-4o")
    parser.add_argument("--input", type=str, required=True, help="The topic for the presentation.")
    parser.add_argument("--slides", type=int, default=3, help="The number of slides to generate.")
    parser.add_argument("--output", type=str, default="presentation.pptx", help="The output file path for the .pptx file.")
    parser.add_argument("--model", type=str, default="gpt-4o", help="The LLM model to use (e.g., gpt-4o).")
    args = parser.parse_args()

    # --- Setup ---
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        stream=sys.stdout
    )

    llm, agent = None, None
    try:
        # --- LLM and Agent Initialization ---
        LLMFactory.ensure_initialized()
        llm = LLMFactory.create_llm(
            args.model,
            llm_config=LLMConfig(temperature=0.1)
        )
        
        agent_tools = [generate_slide_content] # Only this tool is needed for the first step
        
        agent_config = AgentConfig(
            name="ContentGeneratorAgent",
            role="An autonomous presentation content creator.",
            description="An agent that generates structured JSON content for a presentation.",
            llm_instance=llm,
            tools=agent_tools,
            system_prompt="You are an AI assistant that generates structured content for presentations. You are an expert at generating structured presentation content based on a user's topic. Your only task is to call the `generate_slide_content` tool.",
            use_xml_tool_format=True,
        )
        
        agent = AgentFactory().create_agent(config=agent_config)

        # --- Step 1: Generate Slide Content ---
        logger.info("Step 1: Generating slide content...")
        agent.start()
        
        # Create a future to await the result from the tool
        content_future = get_tool_result(agent, "generate_slide_content")
        
        # Send the initial message to trigger the agent
        await agent.post_user_message(AgentInputUserMessage(
            content=f"Create presentation content for topic: '{args.input}' with {args.slides} slides."
        ))

        # Wait for the content generation to complete
        slide_content_json = await content_future
        if not slide_content_json:
            raise Exception("Failed to get slide content from the agent.")
        
        content = json.loads(slide_content_json)
        logger.info("Slide content generated successfully.")

        # --- Step 2: Generate Slide Images ---
        logger.info("Step 2: Generating slide images...")
        image_paths = {}
        
        slides_with_images = [s for s in content.get('slides', []) if s.get('type') == 'content' and s.get('image_prompt')]
        
        image_tasks = []
        for slide_data in slides_with_images:
            slide_num = slide_data.get('number')
            prompt = slide_data.get('image_prompt')
            if slide_num and prompt:
                image_tasks.append(generate_slide_image(prompt, slide_num))
        
        generated_paths = await asyncio.gather(*image_tasks)
        
        for i, slide_data in enumerate(slides_with_images):
            slide_num = slide_data.get('number')
            image_paths[slide_num] = generated_paths[i]
            
        logger.info("Slide images generated successfully.")

        # --- Step 3: Create PowerPoint Presentation ---
        logger.info("Step 3: Creating the final PowerPoint presentation...")
        output_path = Path(args.output).resolve()
        image_paths_json = json.dumps(image_paths)
        
        result_message = await create_powerpoint_presentation(slide_content_json, image_paths_json, str(output_path))
        logger.info(result_message)

    except Exception as e:
        logger.error(f"An error occurred during agent execution: {e}", exc_info=True)
    finally:
        # --- Cleanup ---
        logger.info("Shutting down and cleaning up resources.")
        if agent and agent.is_running:
            agent.shutdown()
        
        AgentThreadPoolManager().shutdown()
        
        if TEMP_DIR.exists():
            shutil.rmtree(TEMP_DIR)
        
        if llm:
            await llm.cleanup()

if __name__ == "__main__":
    asyncio.run(main())