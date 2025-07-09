import asyncio
import logging
import os
import json
import requests
from pathlib import Path
from typing import List, Dict, Any, Optional

# --- Dependency Imports ---
# These libraries need to be installed:
# pip install python-pptx openai
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches

# --- Framework Imports ---
# Assuming these are importable from your project structure
from autobyteus.agent.api import create_agent, send_message_to_agent
from autobyteus.agent.agent import Agent
from autobyteus.agent.context.agent_config import AgentConfig
from autobyteus.agent.workspace.base_workspace import BaseAgentWorkspace
from agent.utils.wait_for_idle import wait_for_agent_to_be_idle
from tools.base_tool import BaseTool
from tools.parameter_schema import ParameterSchema, ParameterDefinition, ParameterType
from tools.registry import default_tool_registry

# --- Setup Logging ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


# ==============================================================================
# 1. DEFINE A CUSTOM WORKSPACE FOR LOCAL FILE HANDLING
# ==============================================================================
class LocalFileWorkspace(BaseAgentWorkspace):
    """A concrete workspace that stores agent artifacts in a local directory."""
    def __init__(self, root_path: str = "agent_workspace"):
        super().__init__()
        self.root_path = Path(root_path).resolve()
        self.root_path.mkdir(parents=True, exist_ok=True)
        logger.info(f"LocalFileWorkspace initialized at: {self.root_path}")

    @classmethod
    def get_type_name(cls) -> str:
        return "LocalFileWorkspace"

    @classmethod
    def get_description(cls) -> str:
        return "A workspace that stores files on the local filesystem."

    @classmethod
    def get_config_schema(cls) -> ParameterSchema:
        schema = ParameterSchema()
        schema.add_parameter(ParameterDefinition(
            name="root_path",
            param_type=ParameterType.STRING,
            description="The local directory path for the workspace.",
            required=False,
            default_value="agent_workspace"
        ))
        return schema

    def get_full_path(self, filename: str) -> Path:
        """Returns the full, absolute path for a file in the workspace."""
        return self.root_path / filename

    def write_file_bytes(self, filename: str, content: bytes) -> str:
        """Writes bytes to a file in the workspace and returns its relative path."""
        full_path = self.get_full_path(filename)
        with open(full_path, "wb") as f:
            f.write(content)
        logger.info(f"Wrote {len(content)} bytes to workspace file: {full_path}")
        return filename  # Return the relative path for the LLM

    def save_presentation(self, pres: Presentation, filename: str) -> str:
        """Saves a python-pptx Presentation object to the workspace."""
        full_path = self.get_full_path(filename)
        pres.save(full_path)
        logger.info(f"Saved presentation to workspace file: {full_path}")
        return filename

# ==============================================================================
# 2. DEFINE THE REQUIRED TOOLS
# ==============================================================================
class ImageGenerationTool(BaseTool):
    """A tool to generate an image from a prompt using DALL-E 3 and save it."""
    TOOL_NAME = "generate_image"

    def __init__(self):
        super().__init__()
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY environment variable is not set. This tool cannot function.")
        self.client = OpenAI(api_key=api_key)

    @classmethod
    def get_name(cls) -> str:
        return cls.TOOL_NAME

    @classmethod
    def get_description(cls) -> str:
        return "Generates an image based on a textual prompt using DALL-E 3 and saves it to the workspace. Returns the relative file path of the saved image."

    @classmethod
    def get_argument_schema(cls) -> Optional[ParameterSchema]:
        schema = ParameterSchema()
        schema.add_parameter(ParameterDefinition(
            name="prompt",
            param_type=ParameterType.STRING,
            description="A detailed textual description of the desired image.",
            required=True
        ))
        schema.add_parameter(ParameterDefinition(
            name="filename",
            param_type=ParameterType.STRING,
            description="The desired filename for the output image (e.g., 'slide_1_image.png').",
            required=True
        ))
        return schema

    async def _execute(self, context: 'AgentContext', prompt: str, filename: str) -> str:
        if not isinstance(context.workspace, LocalFileWorkspace):
            return "Error: This tool requires a LocalFileWorkspace."

        logger.info(f"Generating image with prompt: '{prompt[:50]}...'")
        try:
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            image_url = response.data[0].url
            image_data = requests.get(image_url).content
            
            saved_path = context.workspace.write_file_bytes(filename, image_data)
            return f"Image successfully generated and saved to workspace at: {saved_path}"
        except Exception as e:
            logger.error(f"Failed to generate or save image: {e}", exc_info=True)
            return f"Error: Failed to generate image. Details: {str(e)}"

class PowerpointCreationTool(BaseTool):
    """A tool to create a complete PowerPoint presentation from a list of slides."""
    TOOL_NAME = "create_powerpoint"

    @classmethod
    def get_name(cls) -> str:
        return cls.TOOL_NAME

    @classmethod
    def get_description(cls) -> str:
        return "Creates a .pptx file from a structured list of slides. Each slide must have a title, content, and the path to an image already in the workspace."

    @classmethod
    def get_argument_schema(cls) -> Optional[ParameterSchema]:
        schema = ParameterSchema()
        schema.add_parameter(ParameterDefinition(
            name="filename",
            param_type=ParameterType.STRING,
            description="The final filename for the presentation (e.g., 'presentation.pptx').",
            required=True
        ))
        schema.add_parameter(ParameterDefinition(
            name="slides",
            param_type=ParameterType.STRING,
            description="A JSON string representing a list of slides. Each slide is an object with 'title', 'content', and 'image_path' keys. The 'image_path' must be a relative path within the agent's workspace.",
            required=True
        ))
        return schema

    async def _execute(self, context: 'AgentContext', filename: str, slides: str) -> str:
        if not isinstance(context.workspace, LocalFileWorkspace):
            return "Error: This tool requires a LocalFileWorkspace."

        try:
            slides_data: List[Dict[str, str]] = json.loads(slides)
        except json.JSONDecodeError as e:
            return f"Error: The 'slides' argument was not valid JSON. Details: {e}"

        logger.info(f"Creating presentation '{filename}' with {len(slides_data)} slides.")
        pres = Presentation()
        # Use a layout with Title and Content
        title_content_layout = pres.slide_layouts[5]

        for i, slide_info in enumerate(slides_data):
            slide = pres.slides.add_slide(title_content_layout)
            
            # Set Title and Content
            title = slide.shapes.title
            title.text = slide_info.get("title", f"Slide {i+1}")
            
            content_shape = slide.placeholders[1]
            content_shape.text = slide_info.get("content", "")

            # Add Image
            image_path_str = slide_info.get("image_path")
            if image_path_str:
                try:
                    full_image_path = context.workspace.get_full_path(image_path_str)
                    # Add image to the right side of the slide
                    pres.slides[i].shapes.add_picture(str(full_image_path), Inches(5), Inches(1.5), width=Inches(4.5))
                except FileNotFoundError:
                    logger.warning(f"Image not found at workspace path: {image_path_str}")
                except Exception as e:
                    logger.error(f"Could not add image {image_path_str} to slide {i}: {e}")

        saved_path = context.workspace.save_presentation(pres, filename)
        return f"Presentation successfully created and saved to workspace at: {saved_path}"

# ==============================================================================
# 3. DEFINE THE AGENT CONFIGURATION
# ==============================================================================
POWERPOINT_AGENT_SYSTEM_PROMPT = """You are an expert presentation creator named 'PrezAI'.
Your goal is to take a user's request (a title or a paragraph) and turn it into a complete PowerPoint presentation with text and relevant images.

Follow this sequence of actions strictly:
1.  **Plan the Slides**: First, think about the structure of the presentation. Break down the user's topic into a logical sequence of slides (e.g., Title, Introduction, Key Points, Conclusion).
2.  **Generate Image Prompts**: For each slide in your plan, create a detailed and descriptive prompt suitable for generating a compelling image with a model like DALL-E 3.
3.  **Execute Image Generation**: Use the `generate_image` tool for EACH slide to create the images one by one. You must provide a unique filename for each image (e.g., 'slide_1.png', 'slide_2.png'). Wait for the tool to return the path before proceeding to the next image.
4.  **Assemble the Presentation**: Once ALL images have been generated and you have their file paths, you MUST collect all the slide information (titles, content, and the image file paths returned by the `generate_image` tool) into a single JSON structure.
5.  **Create the PowerPoint File**: Finally, call the `create_powerpoint` tool ONCE. Pass the presentation's desired filename and the complete JSON structure of all slides to the 'slides' argument. The JSON must be passed as a single string.
6.  **Final Response**: After the `create_powerpoint` tool succeeds, inform the user that the presentation has been created and provide the final file path from the tool's output.

Example of the JSON for the 'slides' argument in the `create_powerpoint` tool:
'[{"title": "The Future of AI", "content": "Exploring the next wave of innovation.", "image_path": "slide_1.png"}, {"title": "Key Developments", "content": "- Large Language Models\\n- Generative Art\\n- Autonomous Systems", "image_path": "slide_2.png"}]'

Do not try to write the file yourself or guess file paths. Use the tools provided and the paths they return.
"""

def get_powerpoint_agent_config(workspace: BaseAgentWorkspace, llm_model: str) -> AgentConfig:
    """Creates the AgentConfig for our PowerPoint agent."""
    
    # Register the tools with the default registry so the agent can find them
    default_tool_registry.register_tool(ImageGenerationTool())
    default_tool_registry.register_tool(PowerpointCreationTool())
    
    config = AgentConfig(
        name="PowerPointAgent",
        role="Presentation Creator",
        description="An agent that creates PowerPoint presentations with text and AI-generated images.",
        llm_instance=None,  # The create_agent function will handle this
        system_prompt=POWERPOINT_AGENT_SYSTEM_PROMPT,
        tools=[ImageGenerationTool(), PowerpointCreationTool()],
        workspace=workspace,
        auto_execute_tools=True,
        use_xml_tool_format=False  # Important for models that use JSON for tool calls
    )
    return config

# ==============================================================================
# 4. MAIN RUNNER SCRIPT
# =_============================================================================
async def main():
    """Main function to configure and run the PowerPoint agent."""
    import argparse
    
    parser = argparse.ArgumentParser(description="Run the PowerPoint Generation Agent.")
    parser.add_argument("--topic", type=str, required=True, help="The topic for the presentation.")
    parser.add_argument("--model", type=str, default="gpt-4o", help="The LLM model to use (e.g., 'gpt-4o', 'gpt-4-turbo').")
    parser.add_argument("--workspace", type=str, default="presentations", help="Directory to save generated files.")
    parser.add_argument("--timeout", type=float, default=600.0, help="Timeout in seconds to wait for the agent to finish.")

    args = parser.parse_args()

    # 1. Create and configure the workspace
    workspace = LocalFileWorkspace(root_path=args.workspace)
    logger.info(f"Workspace is ready at {workspace.root_path.resolve()}")

    # 2. Get the agent configuration
    agent_config = get_powerpoint_agent_config(workspace, args.model)

    # 3. Create the agent instance
    logger.info(f"Creating agent with model: {args.model}")
    powerpoint_agent = await create_agent(
        name=agent_config.name,
        role=agent_config.role,
        description=agent_config.description,
        system_prompt=agent_config.system_prompt,
        llm_model=args.model,
        tools=["generate_image", "create_powerpoint"], # Pass tool names for registration lookup
        auto_execute_tools=agent_config.auto_execute_tools,
        use_xml_tool_format=agent_config.use_xml_tool_format,
        # The workspace needs to be set up on the config object for create_agent to use it
    )
    # Manually assign the pre-configured workspace instance
    powerpoint_agent.context.state.workspace = workspace
    
    logger.info(f"Agent '{powerpoint_agent.agent_id}' created successfully.")

    # 4. Start the agent and send the initial message
    try:
        await send_message_to_agent(powerpoint_agent, args.topic)

        # 5. Wait for the agent to complete its task and become idle
        logger.info(f"Agent is processing... Waiting for it to become idle (timeout: {args.timeout}s).")
        await wait_for_agent_to_be_idle(powerpoint_agent, timeout=args.timeout)
        logger.info("Agent has finished its task and is now idle.")

    except asyncio.TimeoutError:
        logger.error(f"Agent did not complete within the {args.timeout}s timeout.")
    except Exception as e:
        logger.error(f"An error occurred while running the agent: {e}", exc_info=True)
    finally:
        # 6. Stop the agent gracefully
        logger.info("Stopping agent...")
        await powerpoint_agent.stop()
        logger.info("Agent stopped.")


if __name__ == "__main__":
    # Ensure you have set the OPENAI_API_KEY environment variable
    if not os.getenv("OPENAI_API_KEY"):
        print("\n" + "="*80)
        print("ERROR: The OPENAI_API_KEY environment variable is not set.")
        print("Please set it to your OpenAI API key to run this agent.")
        print("e.g., export OPENAI_API_KEY='sk-...'")
        print("="*80 + "\n")
    else:
        try:
            asyncio.run(main())
        except KeyboardInterrupt:
            logger.info("Process interrupted by user. Shutting down.")