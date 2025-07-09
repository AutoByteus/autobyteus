#!/usr/bin/env python
# -*- coding: utf-8 -*-

import asyncio
import logging
import argparse
from pathlib import Path
import sys
import os
import json
import tempfile
import shutil
import httpx

# --- Core System Setup ---
# Add the project root to the Python path to ensure all modules are found.
SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

# --- Dependency Imports ---
try:
    from dotenv import load_dotenv
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import openai
except ImportError as e:
    print(f"Error: A required package is missing. Please install it. Details: {e}")
    sys.exit(1)

# --- Autobyteus Framework Imports ---
from autobyteus.agent.context.agent_config import AgentConfig
from autobyteus.agent.factory.agent_factory import AgentFactory
from autobyteus.agent.message.agent_input_user_message import AgentInputUserMessage
from autobyteus.agent.utils.wait_for_idle import wait_for_agent_to_be_idle
from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.llm.utils.llm_config import LLMConfig
from autobyteus.tools.functional_tool import tool
from autobyteus.agent.runtime.agent_thread_pool_manager import AgentThreadPoolManager

# --- Environment and Logging Configuration ---
env_file_path = PROJECT_ROOT / ".env"
if env_file_path.exists():
    print(f"Loading environment variables from: {env_file_path}")
    load_dotenv(dotenv_path=env_file_path)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("PowerPointAgentApp")

# --- Global State and Setup ---
TEMP_DIR = Path(tempfile.mkdtemp(prefix="pptx_agent_"))
logger.info(f"Temporary directory for images created at: {TEMP_DIR}")

# --- Tool Definitions ---

@tool(name="generate_slide_content")
async def generate_slide_content(user_input: str, num_slides: int = 5) -> str:
    """
    Generates structured JSON content for a presentation based on a user's topic.

    Args:
        user_input: The central topic or title for the presentation.
        num_slides: The desired number of content slides, in addition to the title slide.

    Returns:
        A JSON string containing the presentation title and a list of slide objects.
    """
    logger.info(f"Generating slide content for topic: '{user_input}' ({num_slides} slides)")
    llm = LLMFactory.create_llm("gpt-4o", llm_config=LLMConfig(temperature=0.7))
    prompt = f"""
    Based on the user input "{user_input}", generate the content for a {num_slides}-slide presentation.
    The output must be a single, valid JSON object. Do not include any text or formatting outside of the JSON.
    The JSON structure should be:
    {{
      "presentation_title": "A concise, catchy title for the presentation",
      "slides": [
        {{
          "number": 1,
          "type": "title",
          "title": "The main title of the presentation",
          "subtitle": "A brief, engaging subtitle",
          "notes": "Speaker notes for the title slide.",
          "image_prompt": "A simple, professional DALL-E 3 style prompt for a background image for the title slide."
        }},
        {{
          "number": 2,
          "type": "content",
          "title": "Title for Slide 2",
          "bullets": [
            "A key point or bullet for this slide.",
            "Another key point.",
            "A final key point, keep it concise."
          ],
          "notes": "Speaker notes explaining the content of slide 2 in more detail.",
          "image_prompt": "A simple, professional DALL-E 3 style prompt for an image that visually represents the content of this slide."
        }}
      ]
    }}
    Ensure every slide has a unique `image_prompt`. The prompts should be creative and relevant.
    """
    try:
        response = await llm.send_user_message(prompt)
        json_text = response.content.strip().replace("```json", "").replace("```", "").strip()
        json.loads(json_text)  # Validate JSON
        logger.info(f"Successfully generated JSON content for {num_slides + 1} slides.")
        return json_text
    except Exception as e:
        logger.error(f"Failed to generate or parse slide content JSON: {e}", exc_info=True)
        fallback_json = json.dumps({
            "presentation_title": user_input,
            "slides": [{
                "number": 1, "type": "title", "title": user_input,
                "subtitle": "AI-Generated Presentation", "notes": "This is a fallback slide due to an error.",
                "image_prompt": "A generic professional background image for a business presentation."
            }]
        })
        return fallback_json
    finally:
        await llm.cleanup()

@tool(name="generate_slide_image")
async def generate_slide_image(image_prompt: str, slide_number: int) -> str:
    """
    Generates an image using DALL-E 3 based on a text prompt and saves it locally.

    Args:
        image_prompt: A descriptive prompt for the image to be generated.
        slide_number: The number of the slide this image is for, used for naming the file.

    Returns:
        The local file path of the saved image.
    """
    logger.info(f"Generating DALL-E 3 image for slide {slide_number} with prompt: '{image_prompt}'")
    filepath = TEMP_DIR / f"slide_{slide_number}.png"
    try:
        client = openai.AsyncOpenAI()
        response = await client.images.generate(
            model="dall-e-3",
            prompt=image_prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        image_url = response.data[0].url
        if not image_url:
            raise ValueError("DALL-E 3 API did not return an image URL.")

        async with httpx.AsyncClient() as http_client:
            image_response = await http_client.get(image_url, timeout=60.0)
            image_response.raise_for_status()

        with open(filepath, "wb") as f:
            f.write(image_response.content)
        logger.info(f"Image for slide {slide_number} downloaded and saved to: {filepath}")
        return str(filepath.resolve())
    except Exception as e:
        logger.error(f"DALL-E 3 image generation failed for slide {slide_number}: {e}", exc_info=True)
        img = Image.new('RGB', (1024, 1024), (128, 128, 128))
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except IOError:
            font = ImageFont.load_default()
        draw.text((50, 50), f"Slide {slide_number}\nImage Generation Failed", font=font, fill=(255, 255, 255))
        img.save(filepath)
        return str(filepath.resolve())

@tool(name="create_powerpoint_presentation")
async def create_powerpoint_presentation(slide_content_json: str, image_paths_json: str, output_path: str) -> str:
    """
    Creates a PowerPoint .pptx file from structured content and image paths.

    Args:
        slide_content_json: A JSON string containing the presentation title and slide data.
        image_paths_json: A JSON string mapping slide numbers to their local image file paths.
        output_path: The desired file path for the final .pptx presentation.

    Returns:
        A success message with the final path of the created presentation.
    """
    logger.info(f"Assembling PowerPoint presentation at: {output_path}")
    try:
        content = json.loads(slide_content_json)
        image_paths = {int(k): v for k, v in json.loads(image_paths_json).items()}
    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON format for content or image paths. Details: {e}"

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for slide_data in content.get('slides', []):
        slide_num = slide_data.get('number')
        slide_type = slide_data.get('type')
        if slide_type == 'title':
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = slide_data.get('title', 'Presentation Title')
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get('subtitle', '')
        elif slide_type == 'content':
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1.2))
            title_shape.text_frame.text = slide_data.get('title', '')
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            has_image = isinstance(slide_num, int) and slide_num in image_paths and Path(image_paths.get(slide_num, "")).exists()
            
            if slide_data.get('bullets'):
                text_width = Inches(7.5) if has_image else Inches(15)
                body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), text_width, Inches(6.5))
                tf = body_shape.text_frame
                tf.clear()
                for bullet_text in slide_data['bullets']:
                    p = tf.add_paragraph()
                    p.text = str(bullet_text)
                    p.font.size = Pt(28)
                    p.level = 0
            
            if has_image:
                img_path = image_paths[slide_num]
                slide.shapes.add_picture(img_path, Inches(8), Inches(1.5), width=Inches(7.5))

        if slide_data.get('notes'):
            slide.notes_slide.notes_text_frame.text = slide_data.get('notes')

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info(f"Presentation saved successfully to {output_path}")
    return f"Presentation created and saved to {output_path}"

# --- Agent System Prompt ---
SYSTEM_PROMPT = """
You are a highly efficient assistant that creates PowerPoint presentations.
Your goal is to take a user's topic and create a complete .pptx file.
You must follow this workflow precisely:

1.  **Generate Content:** Call the `generate_slide_content` tool using the user's input topic and desired number of slides. This will return a JSON string with all the text, notes, and image prompts for the presentation.

2.  **Generate Images (Loop):**
    -   Parse the JSON output from the previous step.
    -   Iterate through each slide object in the "slides" array.
    -   For **each slide**, call the `generate_slide_image` tool. Use the `image_prompt` and `number` from the current slide object as arguments.
    -   You will get a local file path for each generated image.

3.  **Collect Image Paths:**
    -   After the loop, you must create a single JSON object that maps each slide number (as a string key) to its corresponding image file path (as a string value).
    -   Example format: `{"1": "/path/to/slide_1.png", "2": "/path/to/slide_2.png", ...}`

4.  **Create Presentation:**
    -   Call the `create_powerpoint_presentation` tool.
    -   For the `slide_content_json` argument, pass the complete, original JSON string from Step 1.
    -   For the `image_paths_json` argument, pass the JSON string of the image path mapping you created in Step 3.
    -   For the `output_path` argument, use the file path provided by the user.

Do not ask for confirmation. Execute the entire workflow automatically and sequentially. The final step must be calling `create_powerpoint_presentation`.
"""

# --- Main Execution Block ---
async def main():
    parser = argparse.ArgumentParser(description="PowerPoint Generation Agent using GPT-4o and DALL-E 3")
    parser.add_argument("--input", type=str, required=True, help="The topic for the presentation.")
    parser.add_argument("--slides", type=int, default=3, help="The number of content slides to generate.")
    parser.add_argument("--output", type=str, default="presentation.pptx", help="The output file path for the .pptx file.")
    parser.add_argument("--model", type=str, default="gpt-4o", help="The orchestrator LLM model to use.")
    args = parser.parse_args()

    agent = None
    try:
        LLMFactory.reinitialize()
        
        logger.info("Initializing PowerPoint Generation Agent...")
        agent_llm = LLMFactory.create_llm(
            args.model,
            llm_config=LLMConfig(temperature=0.1)
        )
        
        agent_tools = [generate_slide_content, generate_slide_image, create_powerpoint_presentation]
        
        agent_config = AgentConfig(
            name="PowerPointGPTAgent",
            role="An autonomous presentation creator.",
            description="Generates complete PowerPoint presentations from a topic, including text and images, using a sequence of tools.",
            llm_instance=agent_llm,
            tools=agent_tools,
            system_prompt=SYSTEM_PROMPT,
            auto_execute_tools=True,
        )
        
        agent = AgentFactory().create_agent(config=agent_config)
        
        initial_message_content = f"""
        Please create a presentation based on the following input.
        User Topic: "{args.input}"
        Number of content slides: {args.slides}
        Final output file path: {Path(args.output).resolve()}
        """
        
        agent_input = AgentInputUserMessage(content=initial_message_content)

        agent.start()
        # FIX: Wait briefly for the agent's worker loop to initialize
        await asyncio.sleep(0.1)
        
        await agent.post_user_message(agent_input)

        logger.info("Agent is running... Waiting for the presentation to be completed. This may take a few minutes.")
        await wait_for_agent_to_be_idle(agent, timeout=600.0)

        logger.info("✅ Agent has completed its task.")
        final_output_path = Path(args.output).resolve()
        logger.info(f"Presentation should be available at: {final_output_path}")

        output_dir = final_output_path.parent
        images_dir = output_dir / f"{final_output_path.stem}_images"
        if TEMP_DIR.exists() and any(TEMP_DIR.iterdir()):
            shutil.copytree(TEMP_DIR, images_dir, dirs_exist_ok=True)
            logger.info(f"📁 Supporting images have been copied to: {images_dir}")
        
    except Exception as e:
        logger.error(f"An error occurred during agent execution: {e}", exc_info=True)
    finally:
        logger.info("Shutting down and cleaning up resources.")
        if agent and agent.is_running:
            # FIX: Use the correct stop() method instead of shutdown()
            await agent.stop()
            logger.info("Agent stopped.")
        
        # Shutdown the global thread pool for a clean exit
        AgentThreadPoolManager().shutdown()

        if TEMP_DIR.exists():
            shutil.rmtree(TEMP_DIR, ignore_errors=True)
            logger.info(f"Temporary directory {TEMP_DIR} cleaned up.")

if __name__ == "__main__":
    if not os.getenv("OPENAI_API_KEY"):
         logger.error("The OPENAI_API_KEY environment variable is not set.")
         logger.error("Please set it in your .env file or system environment to use GPT-4o and DALL-E 3.")
         sys.exit(1)
         
    asyncio.run(main())