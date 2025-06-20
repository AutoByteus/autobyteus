#!/usr/bin/env python3
"""
Gemini-Enhanced Google Slides Example
This script demonstrates how to use Gemini's image generation capabilities
to create visually rich slides with AI-generated content.
"""
import asyncio
import sys
import os
from pathlib import Path
import logging
import base64
import uuid
import mimetypes
from io import BytesIO
from PIL import Image
import http.server
import socketserver
import threading
import tempfile
import shutil

# Ensure the autobyteus package is discoverable
SCRIPT_DIR = Path(__file__).resolve().parent
PACKAGE_ROOT = SCRIPT_DIR.parent
sys.path.insert(0, str(PACKAGE_ROOT))

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(name)s - %(message)s')
logger = logging.getLogger("gemini_slides_example")

# Import necessary modules
from autobyteus.llm.llm_factory import LLMFactory
from autobyteus.llm.utils.llm_config import LLMConfig

# Create temporary directory for generated content
TEMP_DIR = tempfile.mkdtemp(prefix="gemini_slides_")
SERVER_PORT = 8000
SERVER_HOST = "localhost"
file_server = None

def start_file_server():
    """Starts a simple HTTP server to serve generated files"""
    global file_server
    
    os.chdir(TEMP_DIR)
    handler = http.server.SimpleHTTPRequestHandler
    file_server = socketserver.TCPServer((SERVER_HOST, SERVER_PORT), handler)
    
    logger.info(f"Starting file server on http://{SERVER_HOST}:{SERVER_PORT}")
    server_thread = threading.Thread(target=file_server.serve_forever)
    server_thread.daemon = True
    server_thread.start()
    
def stop_file_server():
    """Stops the HTTP server"""
    global file_server
    if file_server:
        logger.info("Shutting down file server")
        file_server.shutdown()

def save_content_to_file(content_bytes, file_extension, base_url=None):
    """
    Saves content bytes to a file and returns its URL
    
    Args:
        content_bytes (bytes): The binary content to save
        file_extension (str): The file extension (e.g., 'png', 'mp4')
        base_url (str, optional): Base URL to use instead of localhost
        
    Returns:
        str: URL to the saved file
    """
    filename = f"{uuid.uuid4()}.{file_extension}"
    filepath = os.path.join(TEMP_DIR, filename)
    
    with open(filepath, "wb") as f:
        f.write(content_bytes)
    
    # Use custom base URL if provided, otherwise use the default localhost URL
    if base_url:
        # Make sure the base URL doesn't have a trailing slash
        base_url = base_url.rstrip('/')
        file_url = f"{base_url}/{filename}"
    else:
        file_url = f"http://{SERVER_HOST}:{SERVER_PORT}/{filename}"
        
    logger.info(f"Saved content to {filepath}, accessible at {file_url}")
    return file_url

async def create_gemini_enhanced_presentation(user_email, title="Gemini-Enhanced Presentation", 
                                    dry_run=False, base_url=None):
    """
    Creates a presentation with slides containing Gemini-generated content.
    
    Args:
        user_email (str): The user's Google email address
        title (str): The title for the presentation
        dry_run (bool): If True, generate content but don't create slides
        base_url (str, optional): Base URL for serving images
    """
    from autobyteus.tools.mcp import McpConfigService, McpConnectionManager

    # Set up MCP components
    logger.info("Setting up MCP connection...")
    config_service = McpConfigService()
    conn_manager = McpConnectionManager(config_service=config_service)
    
    # Configure the MCP server
    server_id = "google-slides-mcp-ws"
    mcp_config = {
        server_id: {
            "transport_type": "websocket",
            "uri": "ws://localhost:8765",
            "enabled": True,
            "tool_name_prefix": "gslides",
        }
    }
    config_service.load_configs(mcp_config)
    
    # Initialize the Gemini LLM
    logger.info("Initializing Gemini LLM...")
    gemini = LLMFactory.create_llm(
        model_identifier="gemini-2.5-pro",
        llm_config=LLMConfig(temperature=0.7)
    )
    
    try:
        # Connect to MCP server
        session = await conn_manager.get_session(server_id)
        await session.transport_strategy.connect()
        
        # Step 1: Create a new presentation
        logger.info(f"Creating presentation with title: {title}")
        create_result = await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "create_presentation", 
                "parameters": {
                    "user_google_email": user_email, 
                    "title": title
                }
            }
        )
        
        # Extract presentation ID from the result
        presentation_id = create_result.get("presentationId")
        
        if not presentation_id:
            logger.error("Failed to extract presentation ID")
            return
        
        logger.info(f"Created presentation with ID: {presentation_id}")
        
        # Step 2: Create slides with Gemini-generated content
        if not dry_run:
            await create_title_slide(session, presentation_id, user_email, title)
            await create_ai_generated_image_slide(session, presentation_id, user_email, gemini, base_url)
            await create_diagram_slide(session, presentation_id, user_email, gemini, base_url)
            await create_data_visualization_slide(session, presentation_id, user_email, gemini, base_url)
            await create_concept_illustration_slide(session, presentation_id, user_email, gemini, base_url)
            await create_video_slide(session, presentation_id, user_email, gemini, base_url)
        else:
            # In dry-run mode, just generate the content but don't create slides
            logger.info("Generating content in dry-run mode...")
            # Generate an image
            image_prompt = "A futuristic city with advanced technology, flying vehicles, and sustainable architecture"
            image_bytes_list = await gemini.generate_image(image_prompt, number_of_images=1)
            if image_bytes_list:
                save_content_to_file(image_bytes_list[0], "png", base_url)
            
            # Generate a diagram
            diagram_prompt = "The architecture of an AI-driven education platform"
            diagram_bytes = await gemini.generate_diagram(diagram_prompt, "flowchart")
            if diagram_bytes:
                save_content_to_file(diagram_bytes, "png", base_url)
            
            # Generate a chart
            viz_prompt = "Student performance metrics before and after using AI tutoring"
            viz_bytes = await gemini.generate_diagram(viz_prompt, diagram_type="chart")
            if viz_bytes:
                save_content_to_file(viz_bytes, "png", base_url)
            
            # Generate a concept illustration
            concept_prompt = "The concept of personalized learning with AI"
            concept_bytes_list = await gemini.generate_image(concept_prompt, number_of_images=1)
            if concept_bytes_list:
                save_content_to_file(concept_bytes_list[0], "png", base_url)
            
            # Generate a video
            video_prompt = "A short video showcasing the benefits of digital transformation"
            video_bytes = await gemini.generate_video(video_prompt)
            if video_bytes:
                save_content_to_file(video_bytes, "mp4", base_url)
        
        logger.info(f"Gemini-enhanced presentation created successfully! ID: {presentation_id}")
        print(f"\nPresentation URL: https://docs.google.com/presentation/d/{presentation_id}/edit")
        return presentation_id
        
    except Exception as e:
        logger.error(f"Error creating Gemini-enhanced presentation: {e}", exc_info=True)
    finally:
        # Clean up
        if conn_manager:
            await conn_manager.cleanup()
        if gemini:
            await gemini.cleanup()

async def create_title_slide(session, presentation_id, user_email, title):
    """Creates a title slide with a heading and subtitle"""
    logger.info("Creating title slide...")
    
    requests = [
        # Create a new slide
        {
            "createSlide": {
                "objectId": "titleSlide",
                "insertionIndex": 0,
                "slideLayoutReference": {
                    "predefinedLayout": "TITLE"
                }
            }
        },
        # Add title text
        {
            "createShape": {
                "objectId": "titleTextBox",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "titleSlide",
                    "size": {
                        "width": {"magnitude": 7000000, "unit": "EMU"},
                        "height": {"magnitude": 2000000, "unit": "EMU"}
                    },
                    "transform": {
                        "scaleX": 1,
                        "scaleY": 1,
                        "translateX": 1500000,
                        "translateY": 1500000,
                        "unit": "EMU"
                    }
                }
            }
        },
        {
            "insertText": {
                "objectId": "titleTextBox",
                "text": title
            }
        },
        # Add subtitle text
        {
            "createShape": {
                "objectId": "titleSubtitle",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "titleSlide",
                    "size": {
                        "width": {"magnitude": 6000000, "unit": "EMU"},
                        "height": {"magnitude": 1000000, "unit": "EMU"}
                    },
                    "transform": {
                        "scaleX": 1,
                        "scaleY": 1,
                        "translateX": 1500000,
                        "translateY": 3500000,
                        "unit": "EMU"
                    }
                }
            }
        },
        {
            "insertText": {
                "objectId": "titleSubtitle",
                "text": "Created with Google Slides API and Gemini AI"
            }
        },
        # Add formatting
        {
            "updateTextStyle": {
                "objectId": "titleTextBox",
                "style": {
                    "fontSize": {
                        "magnitude": 40,
                        "unit": "PT"
                    },
                    "foregroundColor": {
                        "opaqueColor": {
                            "rgbColor": {
                                "red": 0.2,
                                "green": 0.2,
                                "blue": 0.7
                            }
                        }
                    },
                    "bold": True
                },
                "textRange": {
                    "type": "ALL"
                },
                "fields": "fontSize,foregroundColor,bold"
            }
        }
    ]
    
    # Apply the batch update
    await session.transport_strategy.rpc_call(
        "tools/call",
        {
            "tool_name": "batch_update_presentation",
            "parameters": {
                "user_google_email": user_email,
                "presentation_id": presentation_id,
                "requests": requests
            }
        }
    )

async def create_ai_generated_image_slide(session, presentation_id, user_email, gemini, base_url=None):
    """Creates a slide with an AI-generated image"""
    logger.info("Creating AI-generated image slide...")
    
    # Step 1: Generate image using Gemini
    prompt = "A futuristic city with flying cars, holographic billboards, and lush green spaces integrated with technology"
    try:
        image_bytes_list = await gemini.generate_image(prompt, number_of_images=1)
        if not image_bytes_list:
            raise ValueError("Failed to generate image")
        
        image_bytes = image_bytes_list[0]
        
        # Save to file and get URL
        image_url = save_content_to_file(image_bytes, "png", base_url)
        
        # Step 2: Create slide with the image
        requests = [
            # Create a new slide
            {
                "createSlide": {
                    "objectId": "imageSlide",
                    "insertionIndex": 1
                }
            },
            # Add title
            {
                "createShape": {
                    "objectId": "imageTitle",
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": "imageSlide",
                        "size": {
                            "width": {"magnitude": 7000000, "unit": "EMU"},
                            "height": {"magnitude": 1000000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 1000000,
                            "unit": "EMU"
                        }
                    }
                }
            },
            {
                "insertText": {
                    "objectId": "imageTitle",
                    "text": "AI-Generated Image: Future City"
                }
            },
            # Insert the image
            {
                "createImage": {
                    "objectId": "aiGeneratedImage",
                    "url": image_url,
                    "elementProperties": {
                        "pageObjectId": "imageSlide",
                        "size": {
                            "width": {"magnitude": 6000000, "unit": "EMU"},
                            "height": {"magnitude": 3375000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 2500000,
                            "unit": "EMU"
                        }
                    }
                }
            }
        ]
        
        # Apply the batch update
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to create image slide: {e}", exc_info=True)
        # Create slide with error message instead
        requests = [
            {"createSlide": {"objectId": "errorSlide", "insertionIndex": 1}},
            {"createShape": {
                "objectId": "errorText",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "errorSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 1000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 3000000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "errorText", "text": f"Failed to generate image: {str(e)}"}}
        ]
        
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )

async def create_diagram_slide(session, presentation_id, user_email, gemini, base_url=None):
    """Creates a slide with an AI-generated diagram"""
    logger.info("Creating diagram slide...")
    
    try:
        # Generate a diagram using Gemini
        diagram_prompt = "A simple software architecture with frontend, backend, and database layers"
        diagram_type = "architecture"
        diagram_bytes = await gemini.generate_diagram(diagram_prompt, diagram_type)
        
        # Save to file and get URL
        diagram_url = save_content_to_file(diagram_bytes, "png", base_url)
        
        requests = [
            # Create a new slide
            {
                "createSlide": {
                    "objectId": "diagramSlide",
                    "insertionIndex": 2
                }
            },
            # Add title
            {
                "createShape": {
                    "objectId": "diagramTitle",
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": "diagramSlide",
                        "size": {
                            "width": {"magnitude": 7000000, "unit": "EMU"},
                            "height": {"magnitude": 1000000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 1000000,
                            "unit": "EMU"
                        }
                    }
                }
            },
            {
                "insertText": {
                    "objectId": "diagramTitle",
                    "text": "AI-Generated Architecture Diagram"
                }
            },
            # Insert the diagram
            {
                "createImage": {
                    "objectId": "aiGeneratedDiagram",
                    "url": diagram_url,
                    "elementProperties": {
                        "pageObjectId": "diagramSlide",
                        "size": {
                            "width": {"magnitude": 6000000, "unit": "EMU"},
                            "height": {"magnitude": 3375000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 2500000,
                            "unit": "EMU"
                        }
                    }
                }
            }
        ]
        
        # Apply the batch update
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to create diagram slide: {e}", exc_info=True)
        # Create slide with fallback text
        requests = [
            {"createSlide": {"objectId": "diagramFallbackSlide", "insertionIndex": 2}},
            {"createShape": {
                "objectId": "diagramFallbackTitle",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "diagramFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 1000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 1000000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "diagramFallbackTitle", "text": "Architecture Diagram (Text Description)"}},
            {"createShape": {
                "objectId": "diagramDescription",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "diagramFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 4000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 2500000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "diagramDescription", "text": 
                "Architecture Components:\n\n" +
                "1. Frontend Layer\n" +
                "   - Web UI (React/Angular)\n" +
                "   - Mobile App\n\n" +
                "2. Backend Layer\n" +
                "   - API Gateway\n" +
                "   - Microservices\n" +
                "   - Authentication Service\n\n" +
                "3. Database Layer\n" +
                "   - SQL Database\n" +
                "   - NoSQL Database\n" +
                "   - Cache\n"
            }}
        ]
        
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )

async def create_data_visualization_slide(session, presentation_id, user_email, gemini, base_url=None):
    """Creates a slide with an AI-generated data visualization"""
    logger.info("Creating data visualization slide...")
    
    try:
        # Generate a chart/visualization using Gemini
        viz_prompt = "A bar chart showing monthly sales data for Q1 and Q2 with clear labels"
        viz_bytes = await gemini.generate_diagram(viz_prompt, diagram_type="chart")
        
        # Save to file and get URL
        viz_url = save_content_to_file(viz_bytes, "png", base_url)
        
        requests = [
            # Create a new slide
            {
                "createSlide": {
                    "objectId": "chartSlide",
                    "insertionIndex": 3
                }
            },
            # Add title
            {
                "createShape": {
                    "objectId": "chartTitle",
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": "chartSlide",
                        "size": {
                            "width": {"magnitude": 7000000, "unit": "EMU"},
                            "height": {"magnitude": 1000000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 1000000,
                            "unit": "EMU"
                        }
                    }
                }
            },
            {
                "insertText": {
                    "objectId": "chartTitle",
                    "text": "AI-Generated Sales Data Visualization"
                }
            },
            # Insert the chart
            {
                "createImage": {
                    "objectId": "aiGeneratedChart",
                    "url": viz_url,
                    "elementProperties": {
                        "pageObjectId": "chartSlide",
                        "size": {
                            "width": {"magnitude": 6000000, "unit": "EMU"},
                            "height": {"magnitude": 3375000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 2500000,
                            "unit": "EMU"
                        }
                    }
                }
            }
        ]
        
        # Apply the batch update
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to create chart slide: {e}", exc_info=True)
        # Create simple text chart as fallback
        requests = [
            {"createSlide": {"objectId": "chartFallbackSlide", "insertionIndex": 3}},
            {"createShape": {
                "objectId": "chartFallbackTitle",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "chartFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 1000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 1000000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "chartFallbackTitle", "text": "Sales Data (Text Format)"}}
        ]
        
        # Create a text-based table
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )
        
        # Add a table instead
        table_requests = [
            {
                "createTable": {
                    "objectId": "salesTable",
                    "rows": 7,
                    "columns": 3,
                    "elementProperties": {
                        "pageObjectId": "chartFallbackSlide",
                        "size": {
                            "width": {"magnitude": 6000000, "unit": "EMU"},
                            "height": {"magnitude": 3000000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 2500000,
                            "unit": "EMU"
                        }
                    }
                }
            }
        ]
        
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": table_requests
                }
            }
        )
        
        # Populate table cells
        months = ["Month", "January", "February", "March", "April", "May", "June"]
        q1_data = ["Q1 Sales", "$10,000", "$12,500", "$15,000", "", "", ""]
        q2_data = ["Q2 Sales", "", "", "", "$17,500", "$20,000", "$22,500"]
        
        cell_requests = []
        
        for i, month in enumerate(months):
            cell_requests.append({
                "insertText": {
                    "objectId": "salesTable",
                    "cellLocation": {
                        "rowIndex": i,
                        "columnIndex": 0
                    },
                    "text": month
                }
            })
            
        for i, data in enumerate(q1_data):
            if data:
                cell_requests.append({
                    "insertText": {
                        "objectId": "salesTable",
                        "cellLocation": {
                            "rowIndex": i,
                            "columnIndex": 1
                        },
                        "text": data
                    }
                })
                
        for i, data in enumerate(q2_data):
            if data:
                cell_requests.append({
                    "insertText": {
                        "objectId": "salesTable",
                        "cellLocation": {
                            "rowIndex": i,
                            "columnIndex": 2
                        },
                        "text": data
                    }
                })
        
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": cell_requests
                }
            }
        )

async def create_concept_illustration_slide(session, presentation_id, user_email, gemini, base_url=None):
    """Creates a slide with an AI-generated concept illustration"""
    logger.info("Creating concept illustration slide...")
    
    try:
        # Generate a concept illustration
        concept_prompt = "A visual metaphor for digital transformation showing a butterfly emerging from binary code"
        concept_bytes_list = await gemini.generate_image(concept_prompt, number_of_images=1)
        concept_bytes = concept_bytes_list[0]
        
        # Save to file and get URL
        concept_url = save_content_to_file(concept_bytes, "png", base_url)
        
        requests = [
            # Create a new slide
            {
                "createSlide": {
                    "objectId": "conceptSlide",
                    "insertionIndex": 4
                }
            },
            # Add title
            {
                "createShape": {
                    "objectId": "conceptTitle",
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": "conceptSlide",
                        "size": {
                            "width": {"magnitude": 7000000, "unit": "EMU"},
                            "height": {"magnitude": 1000000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 1000000,
                            "unit": "EMU"
                        }
                    }
                }
            },
            {
                "insertText": {
                    "objectId": "conceptTitle",
                    "text": "Digital Transformation: A Visual Metaphor"
                }
            },
            # Insert the concept illustration
            {
                "createImage": {
                    "objectId": "conceptImage",
                    "url": concept_url,
                    "elementProperties": {
                        "pageObjectId": "conceptSlide",
                        "size": {
                            "width": {"magnitude": 6000000, "unit": "EMU"},
                            "height": {"magnitude": 3375000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 2500000,
                            "unit": "EMU"
                        }
                    }
                }
            }
        ]
        
        # Apply the batch update
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to create concept illustration slide: {e}", exc_info=True)
        # Create fallback slide
        requests = [
            {"createSlide": {"objectId": "conceptFallbackSlide", "insertionIndex": 4}},
            {"createShape": {
                "objectId": "conceptFallbackTitle",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "conceptFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 1000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 1000000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "conceptFallbackTitle", "text": "Digital Transformation: Key Concepts"}},
            {"createShape": {
                "objectId": "conceptDescription",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "conceptFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 4000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 2500000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "conceptDescription", "text": 
                "Digital Transformation Journey:\n\n" +
                "1. Legacy Systems → Modern Infrastructure\n\n" +
                "2. Manual Processes → Automation\n\n" +
                "3. Static Data → Real-time Analytics\n\n" +
                "4. Isolated Teams → Collaborative Culture\n\n" +
                "5. Product Focus → Customer Experience Focus"
            }}
        ]
        
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )

async def create_video_slide(session, presentation_id, user_email, gemini, base_url=None):
    """Creates a slide with an AI-generated video"""
    logger.info("Creating video slide...")
    
    try:
        # Generate a video using Gemini
        video_prompt = "A short video showcasing the benefits of digital transformation"
        video_bytes = await gemini.generate_video(video_prompt)
        
        # Save video to file and get URL
        video_url = save_content_to_file(video_bytes, "mp4", base_url)
        
        requests = [
            # Create a new slide
            {
                "createSlide": {
                    "objectId": "videoSlide",
                    "insertionIndex": 5
                }
            },
            # Add title
            {
                "createShape": {
                    "objectId": "videoTitle",
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": "videoSlide",
                        "size": {
                            "width": {"magnitude": 7000000, "unit": "EMU"},
                            "height": {"magnitude": 1000000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 1000000,
                            "unit": "EMU"
                        }
                    }
                }
            },
            {
                "insertText": {
                    "objectId": "videoTitle",
                    "text": "AI-Generated Video: Digital Transformation Benefits"
                }
            },
            # Insert the video
            {
                "createVideo": {
                    "objectId": "aiGeneratedVideo",
                    "source": "YOUTUBE",
                    "url": video_url,
                    "elementProperties": {
                        "pageObjectId": "videoSlide",
                        "size": {
                            "width": {"magnitude": 6000000, "unit": "EMU"},
                            "height": {"magnitude": 3375000, "unit": "EMU"}
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 1500000,
                            "translateY": 2500000,
                            "unit": "EMU"
                        }
                    }
                }
            }
        ]
        
        # Apply the batch update
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )
        
    except Exception as e:
        logger.error(f"Failed to create video slide: {e}", exc_info=True)
        # Create fallback slide
        requests = [
            {"createSlide": {"objectId": "videoFallbackSlide", "insertionIndex": 5}},
            {"createShape": {
                "objectId": "videoFallbackTitle",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "videoFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 1000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 1000000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "videoFallbackTitle", "text": "Digital Transformation: Key Benefits"}},
            {"createShape": {
                "objectId": "videoDescription",
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": "videoFallbackSlide",
                    "size": {"width": {"magnitude": 6000000, "unit": "EMU"}, "height": {"magnitude": 4000000, "unit": "EMU"}},
                    "transform": {"scaleX": 1, "scaleY": 1, "translateX": 1500000, "translateY": 2500000, "unit": "EMU"}
                }
            }},
            {"insertText": {"objectId": "videoDescription", "text": 
                "Digital Transformation Benefits:\n\n" +
                "1. Increased Efficiency\n" +
                "2. Enhanced Customer Experience\n" +
                "3. Scalability\n" +
                "4. Cost Reduction\n" +
                "5. Innovation"
            }}
        ]
        
        await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "batch_update_presentation",
                "parameters": {
                    "user_google_email": user_email,
                    "presentation_id": presentation_id,
                    "requests": requests
                }
            }
        )

async def download_presentation_as_pptx(presentation_id, output_file, user_email):
    """
    Downloads a Google Slides presentation as a PowerPoint file using the Google Drive API
    
    Args:
        presentation_id (str): The ID of the Google Slides presentation
        output_file (str): Path where the PowerPoint file will be saved
        user_email (str): The user's Google email address
        
    Returns:
        bool: True if successful, False otherwise
    """
    from autobyteus.tools.mcp import McpConfigService, McpConnectionManager
    
    logger.info(f"Downloading presentation {presentation_id} as PowerPoint...")
    
    # Set up MCP components
    config_service = McpConfigService()
    conn_manager = McpConnectionManager(config_service=config_service)
    
    # Configure the MCP server
    server_id = "google-slides-mcp-ws"
    mcp_config = {
        server_id: {
            "transport_type": "websocket",
            "uri": "ws://localhost:8765",
            "enabled": True,
            "tool_name_prefix": "gslides",
        }
    }
    config_service.load_configs(mcp_config)
    
    try:
        # Connect to MCP server
        session = await conn_manager.get_session(server_id)
        await session.transport_strategy.connect()
        
        # Get Google Drive credentials
        creds_result = await session.transport_strategy.rpc_call(
            "tools/call",
            {
                "tool_name": "get_credentials", 
                "parameters": {
                    "user_google_email": user_email
                }
            }
        )
        
        if not creds_result or "credentials" not in creds_result:
            logger.error("Failed to get Google credentials")
            return await create_placeholder_pptx(presentation_id, output_file, TEMP_DIR)
            
        # Use the Google Drive API directly
        from googleapiclient.discovery import build
        from google.oauth2.credentials import Credentials
        import io
        
        # Create credentials object from the returned credentials
        credentials = Credentials(
            token=creds_result["credentials"]["token"],
            refresh_token=creds_result["credentials"]["refresh_token"],
            token_uri=creds_result["credentials"]["token_uri"],
            client_id=creds_result["credentials"]["client_id"],
            client_secret=creds_result["credentials"]["client_secret"],
            scopes=creds_result["credentials"]["scopes"]
        )
        
        # Build the Drive API client
        drive_service = build('drive', 'v3', credentials=credentials)
        
        # Export the presentation as PowerPoint
        request = drive_service.files().export_media(
            fileId=presentation_id,
            mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
        # Download the file content
        file_content = io.BytesIO()
        downloader = io.BytesIO()
        
        # Execute the request and store the content
        response = request.execute()
        downloader.write(response)
        downloader.seek(0)
        
        # Create output directory if it doesn't exist
        os.makedirs(os.path.dirname(os.path.abspath(output_file)), exist_ok=True)
        
        # Save the PowerPoint file
        with open(output_file, "wb") as f:
            f.write(downloader.read())
            
        logger.info(f"Presentation successfully exported as PowerPoint: {output_file}")
        return True
        
    except Exception as e:
        logger.error(f"Error exporting presentation: {e}", exc_info=True)
        return await create_placeholder_pptx(presentation_id, output_file, TEMP_DIR)
    finally:
        # Clean up
        if conn_manager:
            await conn_manager.cleanup()

async def create_placeholder_pptx(presentation_id, output_file, temp_dir):
    """
    Creates a placeholder PowerPoint file with the generated content
    
    Args:
        presentation_id (str): The ID of the Google Slides presentation
        output_file (str): Path where the PowerPoint file will be saved
        temp_dir (str): Path to the temporary directory with generated content
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import glob
        
        logger.info(f"Creating placeholder PowerPoint with generated content")
        
        # Create a new presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Gemini-Generated Content"
        subtitle.text = "This presentation contains AI-generated content\nView online: https://docs.google.com/presentation/d/" + presentation_id + "/edit"
        
        # Find all generated images in the temp directory
        image_files = glob.glob(os.path.join(temp_dir, "*.png"))
        
        # Add slides with the generated images
        for i, image_file in enumerate(image_files):
            # Add a new slide for each image
            slide_layout = prs.slide_layouts[5]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add a title based on the image index
            if i == 0:
                title = "AI-Generated Image"
            elif i == 1:
                title = "AI-Generated Diagram"
            elif i == 2:
                title = "AI-Generated Data Visualization"
            else:
                title = f"AI-Generated Content {i+1}"
                
            # Add title to the slide
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = title
            title_para.font.bold = True
            title_para.font.size = Pt(32)
            
            # Add the image to the slide
            try:
                slide.shapes.add_picture(image_file, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))
            except Exception as img_error:
                logger.error(f"Error adding image to slide: {img_error}")
                # Add a textbox explaining the error
                txt_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                txt_frame = txt_box.text_frame
                txt_frame.text = f"Image could not be added to PowerPoint.\nOriginal file: {os.path.basename(image_file)}"
        
        # Find video file if any
        video_files = glob.glob(os.path.join(temp_dir, "*.mp4"))
        if video_files:
            # Add a slide for the video
            slide_layout = prs.slide_layouts[5]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title to the slide
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = "AI-Generated Video"
            title_para.font.bold = True
            title_para.font.size = Pt(32)
            
            # Add a note about the video (PowerPoint can't embed videos programmatically with python-pptx)
            txt_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            txt_frame = txt_box.text_frame
            txt_frame.text = f"Video file: {os.path.basename(video_files[0])}\n\nNote: The video file is available in the temporary directory but cannot be embedded directly in PowerPoint using this script."
        
        # Add a final slide with information
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "About This Presentation"
        content.text = "This PowerPoint file contains AI-generated content created with Gemini.\n\n" + \
                      "The original presentation is available online at:\n" + \
                      f"https://docs.google.com/presentation/d/{presentation_id}/edit\n\n" + \
                      "Generated content includes:\n" + \
                      f"• {len(image_files)} images\n" + \
                      f"• {len(video_files)} videos\n\n" + \
                      "Created with the Gemini Slides Example script"
        
        # Create output directory if it doesn't exist
        output_dir = os.path.dirname(os.path.abspath(output_file))
        os.makedirs(output_dir, exist_ok=True)
        
        # Save the PowerPoint file
        prs.save(output_file)
        logger.info(f"Created placeholder PowerPoint file with generated content: {output_file}")
        return True
    except Exception as e:
        logger.error(f"Failed to create placeholder PowerPoint: {e}")
        return False

async def main():
    """Main function to run the example"""
    import argparse
    
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="Create a Google Slides presentation with Gemini-generated content")
    parser.add_argument("email", help="Email address to share the presentation with")
    parser.add_argument("title", nargs="?", default="Gemini-Enhanced Presentation", help="Title of the presentation")
    parser.add_argument("--dry-run", action="store_true", help="Run without actually creating slides (for testing)")
    parser.add_argument("--base-url", default=f"http://{SERVER_HOST}:{SERVER_PORT}", 
                        help=f"Base URL for serving images (default: http://{SERVER_HOST}:{SERVER_PORT})")
    parser.add_argument("--save-ppt", action="store_true", help="Save the presentation as a PowerPoint file locally")
    parser.add_argument("--output-dir", default="./presentations", help="Directory to save the PowerPoint file (default: ./presentations)")
    parser.add_argument("--save-content", action="store_true", help="Save all generated content (images, videos) to the output directory")
    args = parser.parse_args()
    
    # Update base URL for serving files
    base_url = args.base_url
    if base_url and base_url != f"http://{SERVER_HOST}:{SERVER_PORT}":
        logger.info(f"Using custom base URL: {base_url}")
    else:
        base_url = f"http://{SERVER_HOST}:{SERVER_PORT}"
    
    # Start the file server
    start_file_server()
    
    presentation_id = None
    try:
        # If in dry-run mode, just generate content but don't create slides
        if args.dry_run:
            logger.info("Running in dry-run mode - generating content but not creating slides")
            
        presentation_id = await create_gemini_enhanced_presentation(
            user_email=args.email, 
            title=args.title,
            dry_run=args.dry_run,
            base_url=args.base_url
        )
        
        # Create output directory if it doesn't exist
        output_dir = os.path.abspath(args.output_dir)
        os.makedirs(output_dir, exist_ok=True)
        
        # Download the presentation as PowerPoint if requested
        if args.save_ppt and presentation_id:
            # Set output file path
            safe_title = args.title.replace(' ', '_').replace('/', '_').replace('\\', '_')
            output_file = os.path.join(output_dir, f"{safe_title}.pptx")
            
            # Download the presentation
            success = await download_presentation_as_pptx(presentation_id, output_file, args.email)
            
            if success:
                logger.info(f"Presentation saved as: {output_file}")
                print(f"\nPresentation saved as PowerPoint: {output_file}")
            else:
                logger.warning(f"Failed to save presentation as PowerPoint")
        
        # Copy generated content to output directory if requested
        if args.save_content:
            content_dir = os.path.join(output_dir, "content")
            os.makedirs(content_dir, exist_ok=True)
            
            # Copy all generated content from temp directory to output directory
            for file in os.listdir(TEMP_DIR):
                if file.endswith(('.png', '.jpg', '.jpeg', '.mp4', '.gif')):
                    src_file = os.path.join(TEMP_DIR, file)
                    dst_file = os.path.join(content_dir, file)
                    shutil.copy2(src_file, dst_file)
                    logger.info(f"Copied generated content: {dst_file}")
            
            print(f"\nGenerated content saved to: {content_dir}")
    finally:
        # Stop file server
        stop_file_server()

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("\nScript interrupted by user.")
    except Exception as e:
        logger.error(f"Unhandled error: {e}", exc_info=True)
        sys.exit(1) 