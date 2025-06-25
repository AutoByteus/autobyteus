#/Users/aswin/data/auto/autobyteus/examples/gemini_pptx_generator.py
#!/usr/bin/env python3
"""
Gemini-Powered Presentation Workflow Agent
This script runs a multi-agent workflow using the AutoByteUs framework to generate a 
PowerPoint presentation based on a topic and upload it to Google Slides.
"""
import asyncio
import sys
import os
import logging
from pathlib import Path
import uuid
import tempfile
import shutil
import argparse

from autobyteus.tools.parameter_schema import ParameterDefinition, ParameterType, ParameterSchema

# --- Setup Project Path ---
SCRIPT_DIR = Path(__file__).resolve().parent
PACKAGE_ROOT = SCRIPT_DIR.parent
sys.path.insert(0, str(PACKAGE_ROOT))
TEMP_DIR = Path(tempfile.mkdtemp(prefix="gemini_pptx_"))

# --- Configure Logging ---
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(name)s - %(message)s',
    stream=sys.stderr
)
logger = logging.getLogger("gemini_workflow_agent")

# --- Defer Autobyteus Imports ---
try:
    from autobyteus.llm.llm_factory import default_llm_factory
    from autobyteus.llm.utils.llm_config import LLMConfig
    from autobyteus.tools.base_tool import BaseTool
    from autobyteus.agent.context.agent_config import AgentConfig
    from autobyteus.agent.factory.agent_factory import AgentFactory
    from autobyteus.agent.group.agent_group import AgentGroup
    from autobyteus.tools.mcp import McpConfigService, McpConnectionManager, McpSchemaMapper, McpToolRegistrar
    from autobyteus.tools.registry import default_tool_registry
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(f"Failed to import a required library. Please ensure all dependencies are installed. Error: {e}", file=sys.stderr)
    sys.exit(1)

# --- Local Tool for Presentation Generation ---
class CreateFullPresentationTool(BaseTool):
    """A tool that generates a complete PowerPoint file locally on a given topic."""

    @classmethod
    def get_name(cls) -> str:
        return "create_full_presentation_from_topic"

    @classmethod
    def get_description(cls) -> str:
        return "Generates a complete multi-slide .pptx file locally based on a user-provided topic. Returns the path to the saved file."

    @classmethod
    def get_argument_schema(cls) -> ParameterSchema:
        schema = ParameterSchema()
        schema.add_parameter(ParameterDefinition(
            name="topic", param_type=ParameterType.STRING, description="The central topic for the presentation.", required=True
        ))
        return schema

    async def _execute(self, context, topic: str) -> str:
        """The main logic for creating the presentation file."""
        logger.info(f"Tool '{self.get_name()}' starting presentation generation for topic: '{topic}'")
        
        pro_llm = None
        image_llm = None
        try:
            # Instantiate LLMs needed for content generation
            pro_llm = default_llm_factory.create_llm(model_identifier="gemini-1.5-pro-latest", llm_config=LLMConfig(temperature=0.7))
            image_llm = default_llm_factory.create_llm(model_identifier="imagen-3.0-generate-002", llm_config=LLMConfig(temperature=0.8))

            prs = Presentation()
            
            # --- Slide 1: Title ---
            logger.info("Generating title slide content...")
            title_text = f"Presentation on: {topic}"
            subtitle_response = await pro_llm.generate_content(f"Generate a concise, professional subtitle for a presentation about '{topic}'.")
            
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = subtitle_response.text.strip()

            # --- Slide 2: Image ---
            logger.info("Generating image slide content...")
            image_prompt = f"A professional, high-quality photograph representing the core idea of '{topic}'."
            image_bytes_list = await image_llm.generate_image(image_prompt, number_of_images=1)
            
            if image_bytes_list:
                img_path = TEMP_DIR / f"{uuid.uuid4()}.png"
                with open(img_path, "wb") as f:
                    f.write(image_bytes_list[0])
                
                pic_slide_layout = prs.slide_layouts[5] # Title only
                slide = prs.slides.add_slide(pic_slide_layout)
                slide.shapes.title.text = "Visualizing the Concept"
                slide.shapes.add_picture(str(img_path), Inches(1), Inches(1.5), width=Inches(8))
            
            # --- Slide 3: Key Points (Text) ---
            logger.info("Generating key points slide...")
            points_response = await pro_llm.generate_content(f"List 3-5 main bullet points for a presentation about '{topic}'.")
            
            points_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(points_layout)
            slide.shapes.title.text = "Key Discussion Points"
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # Clear default text
            for point in points_response.text.strip().split('\n'):
                p = tf.add_paragraph()
                p.text = point.lstrip('*- ')
                p.level = 0
            
            # --- Finalization ---
            safe_title = topic.replace(" ", "_").replace("/", "")
            output_filename = TEMP_DIR / f"{safe_title}_{uuid.uuid4()}.pptx"
            prs.save(str(output_filename))
            
            logger.info(f"Successfully generated presentation at: {output_filename}")
            return f"Success! Presentation file created at: {output_filename}"

        except Exception as e:
            logger.error(f"Error in '{self.get_name()}' tool: {e}", exc_info=True)
            return f"Error: Failed to generate presentation. Reason: {e}"
        finally:
            # Clean up LLM instances
            if pro_llm and hasattr(pro_llm, 'cleanup'): await pro_llm.cleanup()
            if image_llm and hasattr(image_llm, 'cleanup'): await image_llm.cleanup()


async def main(args):
    """Main function to set up and run the presentation generation workflow."""
    conn_manager = None
    agent_group = None
    try:
        # --- MCP Setup ---
        logger.info("Setting up MCP connection to Google Slides server...")
        config_service = McpConfigService()
        conn_manager = McpConnectionManager(config_service=config_service)
        mcp_config = {"google-slides-mcp": {"transport_type": "websocket", "uri": "ws://localhost:8765", "tool_name_prefix": "gslides"}}
        config_service.load_configs(mcp_config)
        
        registrar = McpToolRegistrar(
            config_service=config_service,
            schema_mapper=McpSchemaMapper(),
            tool_registry=default_tool_registry
        )
        await registrar.discover_and_register_tools()
        all_definitions = default_tool_registry.get_all_definitions()
        remote_tools = [default_tool_registry.create_tool(name) for name in all_definitions if name.startswith("gslides_")]
        if not remote_tools:
            logger.error("Failed to discover any remote Google Slides tools. Ensure the MCP server is running.")
            return

        # --- Agent and Workflow Setup ---
        logger.info("Setting up agent workflow...")
        agent_factory = AgentFactory()

        # Agent 1: PowerPoint-Creator
        creator_config = AgentConfig(
            name="PowerPoint-Creator", role="creator", description="Generates PPTX files from topics.",
            llm_instance=default_llm_factory.create_llm(model_identifier="gemini-1.5-flash-latest"),
            system_prompt="You are a file creator. Your only job is to use the `create_full_presentation_from_topic` tool when asked.",
            tools=[CreateFullPresentationTool()]
        )

        # Agent 2: Slides-Uploader (Coordinator)
        coordinator_prompt = (
            "You are a presentation workflow coordinator. Your goal is to take a user's topic and create a Google Slides presentation.\n"
            "Workflow:\n"
            "1. You will be given a topic and the user's email address.\n"
            "2. Use the `SendMessageTo` tool to ask the 'creator' agent to generate the presentation file. The creator's role name is 'creator'. The message content should be the topic.\n"
            "3. The creator will respond with the local file path of the generated .pptx file.\n"
            "4. Use the `gslides_upload_and_convert_pptx` tool. Provide the `local_file_path` from the creator's response, use the original topic as the `title`, and pass along the `user_google_email`.\n"
            "5. The upload tool will return the final URL of the Google Slides presentation. Output this URL as your final answer to the user, and nothing else."
        )
        coordinator_config = AgentConfig(
            name="Slides-Uploader", role="coordinator", description="Manages the presentation creation and upload workflow.",
            llm_instance=default_llm_factory.create_llm(model_identifier="gemini-1.5-pro-latest"),
            system_prompt=coordinator_prompt,
            tools=remote_tools  # Note: SendMessageTo is added automatically by AgentGroup
        )
        
        # --- Create and Run Agent Group ---
        agent_group = AgentGroup(
            agent_configs=[creator_config, coordinator_config],
            coordinator_config_name="Slides-Uploader"
        )
        
        logger.info(f"Starting agent group for topic: '{args.topic}' and user: '{args.user_email}'")
        await agent_group.start()

        initial_prompt = f"Please create a presentation on the topic '{args.topic}'. The user's email is '{args.user_email}'."
        final_result = await agent_group.process_task_for_coordinator(initial_prompt)

        print("\n--- Workflow Complete ---")
        print(f"Final Result from Coordinator: {final_result}")
        print("-------------------------\n")

    except Exception as e:
        logger.error(f"An unhandled error occurred in the main workflow: {e}", exc_info=True)
    finally:
        # --- Cleanup ---
        if agent_group:
            logger.info("Stopping agent group...")
            await agent_group.stop()
        if conn_manager:
            logger.info("Cleaning up MCP connection...")
            await conn_manager.cleanup()
        if TEMP_DIR.exists():
            logger.info(f"Removing temporary directory: {TEMP_DIR}")
            shutil.rmtree(TEMP_DIR)
        logger.info("Workflow finished.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Run the Gemini-powered presentation generation and upload workflow.")
    parser.add_argument("topic", type=str, help="The topic for the presentation.")
    parser.add_argument("user_email", type=str, help="The user's Google email address for API access and logging.")
    parser.add_argument("--debug", action="store_true", help="Enable debug logging for all components.")
    
    cli_args = parser.parse_args()

    if cli_args.debug:
        logging.getLogger().setLevel(logging.DEBUG)
        logging.getLogger("autobyteus").setLevel(logging.DEBUG)

    # Check for Google credentials in environment
    if not all(os.getenv(var) for var in ["GOOGLE_CLIENT_ID", "GOOGLE_CLIENT_SECRET", "GOOGLE_REFRESH_TOKEN"]):
        logger.critical("Missing Google API credentials in environment variables (GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, GOOGLE_REFRESH_TOKEN).")
        sys.exit(1)
        
    try:
        asyncio.run(main(cli_args))
    except KeyboardInterrupt:
        print("\nWorkflow interrupted by user.")
    except Exception as e:
        logger.error(f"A critical error occurred: {e}", exc_info=True)
        sys.exit(1)